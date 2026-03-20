#!/usr/bin/env python3
"""
Level 3 PPT Analysis Script
Analyzes all Level 3 PPT files and extracts knowledge points mapped to GESP Levels 5-6
"""

import os
import json
import re
from pptx import Presentation
from collections import defaultdict

# Configuration
INPUT_DIR = "/Users/reacher/Downloads/Archive/Level_3/"
OUTPUT_FILE = "/Users/reacher/Developer/gesp/data/analysis/level3-analysis.json"

# Lesson title mapping from cleanup_knowledge.py
TITLE_MAP = {
    1: "数制转换与位运算",
    2: "枚举进阶",
    3: "递推进阶",
    4: "二分进阶",
    5: "贪心进阶",
    6: "阶段复习与测试",
    7: "递归进阶",
    8: "快速排序",
    9: "归并排序",
    10: "栈进阶",
    11: "队列进阶",
    12: "阶段复习与测试",
    13: "链表进阶",
    14: "二叉树应用",
    15: "二叉树进阶",
    16: "深搜进阶",
    17: "广搜进阶",
    18: "阶段复习与测试",
    19: "动态规划基础",
    20: "动态规划-子序列问题",
    21: "动态规划-背包问题",
    22: "动态规划-背包进阶",
    23: "图的基本应用",
    24: "阶段复习与测试",
}

# Review lessons configuration
REVIEW_LESSONS = {6, 12, 18, 24}
STAGES = {
    1: (1, 5),
    2: (7, 11),
    3: (13, 17),
    4: (19, 23),
}

# GESP Level 5-6 Knowledge Point Definitions
GESP_KNOWLEDGE_POINTS = {
    "l5_1": {"name": "初等数论", "keywords": ["素数", "质数", "GCD", "LCM", "最大公约数", "最小公倍数", "同余", "质因数", "分解", "数论"]},
    "l5_2": {"name": "素数筛法", "keywords": ["筛法", "埃氏筛", "线性筛", "欧拉筛", "素数表"]},
    "l5_3": {"name": "分治算法", "keywords": ["分治", "归并排序", "快速排序", "快排", "merge sort", "quick sort", "分而治之"]},
    "l5_4": {"name": "深度优先搜索", "keywords": ["DFS", "深搜", "深度优先", "回溯", "全排列", "递归搜索"]},
    "l5_5": {"name": "广度优先搜索", "keywords": ["BFS", "广搜", "广度优先", "层次遍历", "最短路径", "队列搜索"]},
    "l5_6": {"name": "动态规划基础", "keywords": ["动态规划", "DP", "记忆化搜索", "状态转移", "递推", "最优子结构"]},
    "l5_7": {"name": "简单背包问题", "keywords": ["背包", "01背包", "完全背包", "多重背包", "分组背包", "滚动数组"]},
    "l5_8": {"name": "树的基本概念", "keywords": ["树", "树形结构", "节点", "父子节点", "根节点", "叶子节点", "树的遍历"]},
    "l5_9": {"name": "二叉树", "keywords": ["二叉树", "binary tree", "左右子树", "二叉树遍历", "前序", "中序", "后序"]},
    "l5_10": {"name": "完全二叉树", "keywords": ["完全二叉树", "满二叉树", "二叉树性质", "节点编号"]},
    "l5_11": {"name": "二叉排序树", "keywords": ["二叉排序树", "BST", "二叉搜索树", "平衡二叉树"]},
    "l5_12": {"name": "图的定义与存储", "keywords": ["图", "graph", "邻接矩阵", "邻接表", "有向图", "无向图", "边", "顶点"]},
    "l6_1": {"name": "哈夫曼树", "keywords": ["哈夫曼树", "Huffman", "最优二叉树", "带权路径", "WPL"]},
    "l6_2": {"name": "哈夫曼编码", "keywords": ["哈夫曼编码", "Huffman编码", "前缀编码", "数据压缩"]},
    "l6_3": {"name": "格雷编码", "keywords": ["格雷编码", "Gray code", "循环码", "相邻编码"]},
    "l6_4": {"name": "搜索综合应用", "keywords": ["搜索综合", "DFS应用", "BFS应用", "连通性", "FloodFill", "泛洪"]},
    "l6_5": {"name": "二维动态规划", "keywords": ["二维DP", "二维动态规划", "区间DP", "状态压缩", "多维状态"]},
    "l6_6": {"name": "DP最值优化", "keywords": ["LIS", "LCS", "最长上升子序列", "最长公共子序列", "最值优化"]},
    "l6_7": {"name": "滚动数组优化", "keywords": ["滚动数组", "空间优化", "降维优化", "状态压缩DP"]},
    "l6_8": {"name": "图的遍历", "keywords": ["图遍历", "图的DFS", "图的BFS", "连通分量", "遍历算法"]},
    "l6_9": {"name": "二叉树搜索", "keywords": ["二叉树搜索", "树的查找", "树的应用"]},
    "l6_10": {"name": "面向对象", "keywords": ["面向对象", "OOP", "类", "class", "封装", "继承", "多态", "对象"]},
}

TOPIC_KEYWORDS = {
    "进制转换": ["l5_1"],
    "位运算": ["l5_1"],
    "枚举": ["l5_3"],
    "递推": ["l5_6"],
    "二分": ["l5_4"],
    "贪心": ["l5_3"],
    "递归": ["l5_4"],
    "排序": ["l5_3"],
    "栈": ["l6_4", "l5_5"],
    "队列": ["l5_5", "l6_4"],
    "链表": ["l5_1"],
    "树": ["l5_8", "l5_9", "l6_1", "l6_2", "l6_3", "l6_4"],
    "二叉树": ["l5_9", "l6_1", "l6_2", "l6_3", "l6_4"],
    "深搜": ["l5_4"],
    "广搜": ["l5_5"],
    "动态规划": ["l5_6", "l5_7", "l6_5", "l6_6", "l6_7"],
    "背包": ["l5_7", "l6_5", "l6_7"],
    "图": ["l5_12", "l6_8"],
}

def extract_text_from_slide(slide):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip())
        if hasattr(shape, "table"):
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        texts.append(cell.text.strip())
    return " ".join(texts)

def identify_knowledge_points(text, lesson_num):
    points = set()
    text_lower = text.lower()
    
    for point_id, point_info in GESP_KNOWLEDGE_POINTS.items():
        for keyword in point_info["keywords"]:
            if keyword.lower() in text_lower or keyword in text:
                points.add(point_id)
                break
    
    for topic, point_ids in TOPIC_KEYWORDS.items():
        if topic in text:
            for pid in point_ids:
                points.add(pid)
    
    lesson_specific_mappings = {
        1: ["l5_1"],
        2: ["l5_3"],
        3: ["l5_6"],
        4: ["l5_4"],
        5: ["l5_3"],
        7: ["l5_4"],
        8: ["l5_3"],
        9: ["l5_3"],
        10: ["l6_4", "l5_5"],
        11: ["l5_5", "l6_4"],
        13: ["l5_1"],
        14: ["l5_9", "l6_2"],
        15: ["l5_9", "l5_8"],
        16: ["l5_4"],
        17: ["l5_5"],
        19: ["l5_6", "l5_7"],
        20: ["l5_6", "l6_6"],
        21: ["l5_7"],
        22: ["l5_7", "l6_7"],
        23: ["l5_12", "l6_8"],
    }
    
    if lesson_num in lesson_specific_mappings and lesson_num not in REVIEW_LESSONS:
        for pid in lesson_specific_mappings[lesson_num]:
            points.add(pid)
    
    return list(points)

def get_lesson_number_from_filename(filename):
    match = re.match(r'L3-(\d+)', filename, re.IGNORECASE)
    if match:
        return int(match.group(1))
    return None

def analyze_ppt(filepath, lesson_num):
    try:
        prs = Presentation(filepath)
    except Exception as e:
        print(f"Error reading {filepath}: {e}")
        return None
    
    pages = []
    all_knowledge_points = set()
    
    for i, slide in enumerate(prs.slides, 1):
        content = extract_text_from_slide(slide)
        if not content:
            continue
        
        content_summary = content[:300] + "..." if len(content) > 300 else content
        points = identify_knowledge_points(content, lesson_num)
        
        page_data = {
            "pageNum": i,
            "content": content_summary,
            "knowledgePoints": points
        }
        pages.append(page_data)
        
        for p in points:
            all_knowledge_points.add(p)
    
    gesp_levels = set()
    for point in all_knowledge_points:
        if point.startswith("l5_"):
            gesp_levels.add(5)
        elif point.startswith("l6_"):
            gesp_levels.add(6)
    
    return {
        "pages": pages,
        "allKnowledgePoints": sorted(list(all_knowledge_points)),
        "gespLevelsCovered": sorted(list(gesp_levels)),
        "totalPages": len(prs.slides)
    }

def main():
    print("=" * 60)
    print("Level 3 PPT Analysis - GESP Levels 5-6")
    print("=" * 60)
    
    ppt_files = []
    for f in os.listdir(INPUT_DIR):
        if f.lower().endswith('.pptx') and f.lower().startswith('l3-'):
            lesson_num = get_lesson_number_from_filename(f)
            if lesson_num:
                ppt_files.append((lesson_num, f))
    
    ppt_files.sort()
    print(f"Found {len(ppt_files)} PPT files to analyze")
    print()
    
    lessons = []
    all_lessons_data = {}
    
    for lesson_num, filename in ppt_files:
        filepath = os.path.join(INPUT_DIR, filename)
        title = TITLE_MAP.get(lesson_num, f"Lesson {lesson_num}")
        
        print(f"Analyzing L3-{lesson_num:02d}: {title}...", end=" ")
        
        analysis = analyze_ppt(filepath, lesson_num)
        if analysis is None:
            print("FAILED")
            continue
        
        is_review = lesson_num in REVIEW_LESSONS
        
        lesson_data = {
            "lessonNum": lesson_num,
            "title": title,
            "fileName": filename,
            "totalPages": analysis["totalPages"],
            "isReview": is_review,
            "stage": 1 if lesson_num <= 6 else 2 if lesson_num <= 12 else 3 if lesson_num <= 18 else 4,
            "pages": analysis["pages"],
            "allKnowledgePoints": analysis["allKnowledgePoints"],
            "gespLevelsCovered": analysis["gespLevelsCovered"],
            "reviewOfLessons": None
        }
        
        all_lessons_data[lesson_num] = lesson_data
        print(f"OK ({analysis['totalPages']} slides, {len(analysis['allKnowledgePoints'])} points)")
        lessons.append(lesson_data)
    
    # Process review lessons
    print()
    print("Processing review lessons...")
    
    for lesson_num in REVIEW_LESSONS:
        if lesson_num not in all_lessons_data:
            continue
        
        review_range = None
        for stage, (start, end) in STAGES.items():
            if lesson_num == 6 and stage == 1:
                review_range = (start, end)
            elif lesson_num == 12 and stage == 2:
                review_range = (start, end)
            elif lesson_num == 18 and stage == 3:
                review_range = (start, end)
            elif lesson_num == 24 and stage == 4:
                review_range = (start, end)
        
        if not review_range:
            continue
        
        reviewed_points = set()
        for ln in range(review_range[0], review_range[1] + 1):
            if ln in all_lessons_data:
                reviewed_points.update(all_lessons_data[ln]["allKnowledgePoints"])
        
        all_lessons_data[lesson_num]["allKnowledgePoints"] = sorted(list(reviewed_points))
        all_lessons_data[lesson_num]["reviewOfLessons"] = list(range(review_range[0], review_range[1] + 1))
        
        gesp_levels = set()
        for point in reviewed_points:
            if point.startswith("l5_"):
                gesp_levels.add(5)
            elif point.startswith("l6_"):
                gesp_levels.add(6)
        all_lessons_data[lesson_num]["gespLevelsCovered"] = sorted(list(gesp_levels))
        
        print(f"  Lesson {lesson_num}: Aggregated {len(reviewed_points)} knowledge points from lessons {review_range[0]}-{review_range[1]}")
    
    lessons = [all_lessons_data[ln] for ln in sorted(all_lessons_data.keys())]
    
    all_points = set()
    gesp5_points = set()
    gesp6_points = set()
    
    for lesson in lessons:
        for point in lesson["allKnowledgePoints"]:
            all_points.add(point)
            if point.startswith("l5_"):
                gesp5_points.add(point)
            elif point.startswith("l6_"):
                gesp6_points.add(point)
    
    total_gesp5 = len([k for k in GESP_KNOWLEDGE_POINTS.keys() if k.startswith("l5_")])
    total_gesp6 = len([k for k in GESP_KNOWLEDGE_POINTS.keys() if k.startswith("l6_")])
    
    result = {
        "level": 3,
        "totalLessons": len(lessons),
        "gespTargetLevels": [5, 6],
        "lessons": lessons,
        "summary": {
            "totalKnowledgePoints": len(all_points),
            "gesp5Coverage": round(len(gesp5_points) / total_gesp5, 2) if total_gesp5 > 0 else 0,
            "gesp6Coverage": round(len(gesp6_points) / total_gesp6, 2) if total_gesp6 > 0 else 0,
            "gesp5PointsCovered": sorted(list(gesp5_points)),
            "gesp6PointsCovered": sorted(list(gesp6_points)),
            "gesp5TotalAvailable": total_gesp5,
            "gesp6TotalAvailable": total_gesp6,
        }
    }
    
    with open(OUTPUT_FILE, 'w', encoding='utf-8') as f:
        json.dump(result, f, ensure_ascii=False, indent=2)
    
    print()
    print("=" * 60)
    print("Analysis Complete!")
    print("=" * 60)
    print(f"Output file: {OUTPUT_FILE}")
    print(f"Total lessons analyzed: {len(lessons)}")
    print(f"Total unique knowledge points: {len(all_points)}")
    print(f"GESP Level 5 coverage: {len(gesp5_points)}/{total_gesp5} ({result['summary']['gesp5Coverage']:.0%})")
    print(f"GESP Level 6 coverage: {len(gesp6_points)}/{total_gesp6} ({result['summary']['gesp6Coverage']:.0%})")

if __name__ == "__main__":
    main()
