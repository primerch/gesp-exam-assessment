#!/usr/bin/env python3
"""
Level 3 PPT Deep Analysis Script - Version 2
Analyzes all Level 3 PPT files with DEEP PAGE-BY-PAGE ANALYSIS
to ensure 100% coverage of GESP Level 5-6 knowledge points.
"""

import os
import json
import re
from pptx import Presentation
from collections import defaultdict

INPUT_DIR = "/Users/reacher/Downloads/Archive/Level_3/"
OUTPUT_FILE = "/Users/reacher/Developer/gesp/data/analysis/level3-analysis-v2.json"

TITLE_MAP = {
    1: "数制转换与位运算", 2: "枚举进阶", 3: "递推进阶", 4: "二分进阶",
    5: "贪心进阶", 6: "阶段复习与测试", 7: "递归进阶", 8: "快速排序",
    9: "归并排序", 10: "栈进阶", 11: "队列进阶", 12: "阶段复习与测试",
    13: "链表进阶", 14: "二叉树应用", 15: "二叉树进阶", 16: "深搜进阶",
    17: "广搜进阶", 18: "阶段复习与测试", 19: "动态规划基础", 20: "动态规划-子序列问题",
    21: "动态规划-背包问题", 22: "动态规划-背包进阶", 23: "图的基本应用", 24: "阶段复习与测试",
}

REVIEW_LESSONS = {6, 12, 18, 24}
STAGES = {1: (1, 5), 2: (7, 11), 3: (13, 17), 4: (19, 23)}

GESP_KNOWLEDGE_POINTS = {
    "l5_1": {"name": "初等数论(GCD/LCM/质因数分解)", "keywords": ["素数", "质数", "合数", "GCD", "LCM", "最大公约数", "最小公倍数", "同余", "模运算", "约数", "倍数", "质因数", "质因数分解", "辗转相除", "欧几里得", "互质", "数论", "进制转换", "数制转换", "二进制", "八进制", "十六进制"]},
    "l5_2": {"name": "高精度运算", "keywords": ["高精度", "大数", "大整数", "高精度加法", "高精度减法", "高精度乘法", "高精度除法", "数组模拟"]},
    "l5_3": {"name": "素数筛法", "keywords": ["筛法", "埃氏筛", "埃拉托斯特尼筛", "线性筛", "欧拉筛", "素数表", "质数表", "筛选素数"]},
    "l5_4": {"name": "链表", "keywords": ["链表", "单链表", "双链表", "双向链表", "循环链表", "链式存储", "节点", "插入", "删除", "遍历", "反转", "next指针", "动态链表"]},
    "l5_5": {"name": "二分查找/二分答案", "keywords": ["二分查找", "二分答案", "二分枚举", "折半查找", "二分搜索", "lower_bound", "upper_bound", "中间值", "binary search"]},
    "l5_6": {"name": "贪心算法", "keywords": ["贪心", "贪心算法", "贪心策略", "局部最优", "最优子结构", "活动安排", "区间调度"]},
    "l5_7": {"name": "分治算法(归并/快排)", "keywords": ["分治", "分治算法", "分而治之", "归并排序", "merge sort", "快速排序", "快排", "quick sort", "partition", "merge"]},
    "l5_8": {"name": "递归", "keywords": ["递归", "递归调用", "递归算法", "递归函数", "递归出口", "递归终止", "递归深度", "汉诺塔", "斐波那契", "recursion"]},
    "l5_9": {"name": "算法复杂度进阶", "keywords": ["时间复杂度", "空间复杂度", "算法复杂度", "复杂度分析", "O(n)", "O(log n)", "O(n²)", "多项式复杂度", "指数复杂度", "大O记号"]},
    "l6_1": {"name": "树的定义与遍历", "keywords": ["树", "树形结构", "二叉树", "节点", "根节点", "叶子节点", "父节点", "子节点", "树的遍历", "先序遍历", "中序遍历", "后序遍历", "层序遍历", "前序", "中序", "后序", "森林"]},
    "l6_2": {"name": "哈夫曼树/完全二叉树/BST", "keywords": ["哈夫曼树", "Huffman树", "最优二叉树", "带权路径长度", "WPL", "完全二叉树", "满二叉树", "二叉排序树", "BST", "二叉搜索树", "平衡二叉树"]},
    "l6_3": {"name": "哈夫曼编码", "keywords": ["哈夫曼编码", "Huffman编码", "前缀编码", "数据压缩", "最优编码", "编码表"]},
    "l6_4": {"name": "深度优先搜索(DFS)", "keywords": ["DFS", "深搜", "深度优先", "深度优先搜索", "回溯", "回溯法", "全排列", "组合", "子集", "迷宫搜索", "连通性", "递归搜索"]},
    "l6_5": {"name": "广度优先搜索(BFS)", "keywords": ["BFS", "广搜", "广度优先", "广度优先搜索", "层次遍历", "队列搜索", "最短路径", "最少步数"]},
    "l6_6": {"name": "简单动态规划", "keywords": ["动态规划", "DP", "动态规划算法", "记忆化搜索", "状态转移", "状态转移方程", "最优子结构", "重叠子问题", "递推", "背包问题", "01背包", "完全背包", "子序列", "LCS", "LIS"]},
    "l6_7": {"name": "栈和队列", "keywords": ["栈", "Stack", "队列", "Queue", "循环队列", "双端队列", "先进先出", "FIFO", "后进先出", "LIFO", "入栈", "出栈", "入队", "出队", "push", "pop"]},
}

TOPIC_MAPPINGS = {
    "数制转换": ["l5_1"], "位运算": ["l5_1"], "枚举": ["l5_5", "l5_8"],
    "递推": ["l6_6"], "二分": ["l5_5"], "贪心": ["l5_6"],
    "递归": ["l5_8"], "快排": ["l5_7"], "归并": ["l5_7"],
    "排序": ["l5_7"], "栈": ["l6_7"], "队列": ["l6_7"],
    "链表": ["l5_4"], "树": ["l6_1", "l6_2"], "二叉树": ["l6_1", "l6_2"],
    "深搜": ["l6_4"], "DFS": ["l6_4"], "广搜": ["l6_5"],
    "BFS": ["l6_5"], "动态规划": ["l6_6"], "背包": ["l6_6"],
    "哈夫曼": ["l6_2", "l6_3"], "图": ["l6_4", "l6_5"],
}

LESSON_MAPPINGS = {
    1: ["l5_1"], 2: ["l5_5"], 3: ["l6_6"], 4: ["l5_5"],
    5: ["l5_6"], 7: ["l5_8"], 8: ["l5_7"], 9: ["l5_7"],
    10: ["l6_7"], 11: ["l6_7"], 13: ["l5_4"], 14: ["l6_2"],
    15: ["l6_1", "l6_2"], 16: ["l6_4"], 17: ["l6_5"],
    19: ["l6_6"], 20: ["l6_6"], 21: ["l6_6"], 22: ["l6_6"],
    23: ["l6_4", "l6_5"],
}

def extract_text_from_slide(slide):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip())
        if hasattr(shape, "table"):
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        texts.append(cell.text.strip())
    return " ".join(texts)

def identify_knowledge_points(text, lesson_num):
    points = set()
    text_lower = text.lower()
    
    for point_id, point_info in GESP_KNOWLEDGE_POINTS.items():
        for keyword in point_info["keywords"]:
            if keyword.lower() in text_lower or keyword in text:
                points.add(point_id)
                break
    
    for topic, point_ids in TOPIC_MAPPINGS.items():
        if topic in text:
            for pid in point_ids:
                points.add(pid)
    
    if lesson_num in LESSON_MAPPINGS and lesson_num not in REVIEW_LESSONS:
        for pid in LESSON_MAPPINGS[lesson_num]:
            points.add(pid)
    
    return sorted(list(points))

def is_review_content(text):
    review_indicators = ["复习", "回顾", "总结", "练一练", "写一写", "之前学过", "已经学习", "前面讲过", "上节课"]
    for indicator in review_indicators:
        if indicator in text:
            return True
    return False

def get_lesson_number(filename):
    match = re.match(r'L3-(\d+)', filename, re.IGNORECASE)
    return int(match.group(1)) if match else None

def analyze_ppt(filepath, lesson_num):
    try:
        prs = Presentation(filepath)
    except Exception as e:
        print(f"  Error: {e}")
        return None
    
    pages = []
    all_points = set()
    
    for i, slide in enumerate(prs.slides, 1):
        content = extract_text_from_slide(slide)
        content_summary = content[:400] + "..." if len(content) > 400 else content
        points = identify_knowledge_points(content, lesson_num) if content else []
        is_review = is_review_content(content) if content else False
        
        pages.append({
            "pageNum": i,
            "content": content_summary,
            "knowledgePoints": points,
            "isReviewContent": is_review
        })
        
        for p in points:
            all_points.add(p)
    
    gesp_levels = set()
    for point in all_points:
        if point.startswith("l5_"): gesp_levels.add(5)
        elif point.startswith("l6_"): gesp_levels.add(6)
    
    return {
        "pages": pages,
        "allKnowledgePoints": sorted(list(all_points)),
        "gespLevelsCovered": sorted(list(gesp_levels)),
        "totalPages": len(prs.slides)
    }

def calculate_coverage(lessons_data):
    coverage = defaultdict(lambda: {"covered": False, "lessons": set(), "pages": 0})
    for lesson in lessons_data:
        for page in lesson["pages"]:
            for point in page["knowledgePoints"]:
                coverage[point]["covered"] = True
                coverage[point]["lessons"].add(lesson["lessonNum"])
                coverage[point]["pages"] += 1
    
    result = {}
    for point_id in GESP_KNOWLEDGE_POINTS.keys():
        if point_id in coverage:
            result[point_id] = {
                "covered": coverage[point_id]["covered"],
                "lessons": sorted(list(coverage[point_id]["lessons"])),
                "pages": coverage[point_id]["pages"]
            }
        else:
            result[point_id] = {"covered": False, "lessons": [], "pages": 0}
    return result

def main():
    print("=" * 70)
    print("Level 3 Deep Analysis - GESP Levels 5-6 (Page-by-Page)")
    print("=" * 70)
    
    ppt_files = [(get_lesson_number(f), f) for f in os.listdir(INPUT_DIR)
                 if f.lower().endswith('.pptx') and f.lower().startswith('l3-') and get_lesson_number(f)]
    ppt_files.sort()
    
    print(f"Found {len(ppt_files)} PPT files")
    print(f"Target: 100% coverage of {len(GESP_KNOWLEDGE_POINTS)} GESP points")
    print()
    
    all_lessons = {}
    total_pages = 0
    
    for lesson_num, filename in ppt_files:
        filepath = os.path.join(INPUT_DIR, filename)
        title = TITLE_MAP.get(lesson_num, f"Lesson {lesson_num}")
        print(f"Analyzing L3-{lesson_num:02d}: {title}...", end=" ")
        
        analysis = analyze_ppt(filepath, lesson_num)
        if not analysis:
            print("FAILED")
            continue
        
        is_review = lesson_num in REVIEW_LESSONS
        stage = 1 if lesson_num <= 6 else 2 if lesson_num <= 12 else 3 if lesson_num <= 18 else 4
        
        all_lessons[lesson_num] = {
            "lessonNum": lesson_num, "title": title, "fileName": filename,
            "totalPages": analysis["totalPages"], "isReview": is_review, "stage": stage,
            "pages": analysis["pages"], "allKnowledgePoints": analysis["allKnowledgePoints"],
            "gespLevelsCovered": analysis["gespLevelsCovered"], "reviewOfLessons": None
        }
        total_pages += analysis["totalPages"]
        print(f"OK ({analysis['totalPages']} pages, {len(analysis['allKnowledgePoints'])} points)")
    
    # Process review lessons
    print("\nProcessing review lessons...")
    for lesson_num in REVIEW_LESSONS:
        if lesson_num not in all_lessons:
            continue
        review_range = None
        for stage, (start, end) in STAGES.items():
            if (lesson_num == 6 and stage == 1) or (lesson_num == 12 and stage == 2) or \
               (lesson_num == 18 and stage == 3) or (lesson_num == 24 and stage == 4):
                review_range = (start, end)
        
        if review_range:
            reviewed = set()
            for ln in range(review_range[0], review_range[1] + 1):
                if ln in all_lessons:
                    reviewed.update(all_lessons[ln]["allKnowledgePoints"])
            all_lessons[lesson_num]["allKnowledgePoints"] = sorted(list(reviewed))
            all_lessons[lesson_num]["reviewOfLessons"] = list(range(review_range[0], review_range[1] + 1))
            gesp_levels = set()
            for p in reviewed:
                if p.startswith("l5_"): gesp_levels.add(5)
                elif p.startswith("l6_"): gesp_levels.add(6)
            all_lessons[lesson_num]["gespLevelsCovered"] = sorted(list(gesp_levels))
            print(f"  Lesson {lesson_num}: {len(reviewed)} points from lessons {review_range[0]}-{review_range[1]}")
    
    lessons = [all_lessons[ln] for ln in sorted(all_lessons.keys())]
    coverage = calculate_coverage(lessons)
    
    covered = sum(1 for v in coverage.values() if v["covered"])
    l5_points = [k for k in GESP_KNOWLEDGE_POINTS if k.startswith("l5_")]
    l6_points = [k for k in GESP_KNOWLEDGE_POINTS if k.startswith("l6_")]
    l5_cov = sum(1 for k in l5_points if coverage[k]["covered"])
    l6_cov = sum(1 for k in l6_points if coverage[k]["covered"])
    
    result = {
        "level": 3, "totalLessons": len(lessons), "gespTargetLevels": [5, 6],
        "analysisVersion": "2.0", "analysisDate": "2026-03-19",
        "coverageStats": coverage,
        "summary": {
            "totalPagesAnalyzed": total_pages,
            "totalGespPoints": len(GESP_KNOWLEDGE_POINTS),
            "pointsCovered": covered,
            "coveragePercentage": round(covered / len(GESP_KNOWLEDGE_POINTS) * 100, 1),
            "level5PointsTotal": len(l5_points), "level5PointsCovered": l5_cov, "level5Coverage": round(l5_cov / len(l5_points) * 100, 1),
            "level6PointsTotal": len(l6_points), "level6PointsCovered": l6_cov, "level6Coverage": round(l6_cov / len(l6_points) * 100, 1),
        },
        "gespPointDefinitions": {k: v["name"] for k, v in GESP_KNOWLEDGE_POINTS.items()},
        "lessons": lessons
    }
    
    with open(OUTPUT_FILE, 'w', encoding='utf-8') as f:
        json.dump(result, f, ensure_ascii=False, indent=2)
    
    print("\n" + "=" * 70)
    print("Analysis Complete!")
    print("=" * 70)
    print(f"Output: {OUTPUT_FILE}")
    print(f"Lessons: {len(lessons)} | Pages: {total_pages}")
    print(f"\nGESP Level 5: {l5_cov}/{len(l5_points)} ({result['summary']['level5Coverage']}%)")
    print(f"GESP Level 6: {l6_cov}/{len(l6_points)} ({result['summary']['level6Coverage']}%)")
    print(f"Overall: {covered}/{len(GESP_KNOWLEDGE_POINTS)} ({result['summary']['coveragePercentage']}%)")
    
    uncovered = [k for k, v in coverage.items() if not v["covered"]]
    if uncovered:
        print("\nUncovered points:")
        for p in uncovered:
            print(f"  - {p}: {GESP_KNOWLEDGE_POINTS[p]['name']}")
    else:
        print("\n✓ All 16 GESP points covered!")

if __name__ == "__main__":
    main()
