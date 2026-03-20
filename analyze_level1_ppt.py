#!/usr/bin/env python3
"""
Level 1 PPT Analysis Script
Analyzes all 24 PPT files and extracts knowledge points from each lesson.
"""

import json
import os
import re
from pptx import Presentation

PPT_DIR = "/Users/reacher/Downloads/Archive/Level_1/"
OUTPUT_FILE = "/Users/reacher/Developer/gesp/data/analysis/level1-analysis.json"

KNOWLEDGE_POINT_MAP = {
    "l1_1": {"name": "计算机基础知识", "keywords": ["计算机", "CPU", "内存", "存储", "硬件", "软件", "操作系统"], "gespLevel": 1},
    "l1_2": {"name": "C++程序框架", "keywords": ["#include", "using namespace", "main函数", "return 0", "头文件"], "gespLevel": 1},
    "l1_3": {"name": "输入输出语句", "keywords": ["cin", "cout", "scanf", "printf", "输入", "输出"], "gespLevel": 1},
    "l1_4": {"name": "变量定义与使用", "keywords": ["变量", "变量定义", "赋值", "初始化", "命名规则"], "gespLevel": 1},
    "l1_5": {"name": "基本数据类型", "keywords": ["int", "long long", "float", "double", "char", "bool", "数据类型"], "gespLevel": 1},
    "l1_6": {"name": "算术运算", "keywords": ["算术运算", "加减乘除", "取余", "模运算", "++", "--"], "gespLevel": 1},
    "l1_7": {"name": "逻辑运算", "keywords": ["逻辑", "&&", "||", "!", "逻辑与", "逻辑或", "逻辑非"], "gespLevel": 1},
    "l1_8": {"name": "关系运算", "keywords": ["关系运算", "==", "!=", "等于", "不等于", "大于", "小于"], "gespLevel": 1},
    "l1_9": {"name": "三目运算", "keywords": ["三目运算", "?:", "条件运算符"], "gespLevel": 1},
    "l1_10": {"name": "顺序结构", "keywords": ["顺序结构", "程序流程"], "gespLevel": 1},
    "l1_11": {"name": "分支结构", "keywords": ["if", "else", "switch", "分支", "条件", "判断"], "gespLevel": 1},
    "l1_12": {"name": "循环结构", "keywords": ["for", "while", "do-while", "循环", "遍历"], "gespLevel": 1},
    "l1_13": {"name": "数组基础", "keywords": ["数组", "一维数组", "arr[", "数组初始化"], "gespLevel": 1},
    "l1_14": {"name": "break和continue", "keywords": ["break", "continue", "跳出", "跳过"], "gespLevel": 1},
    "l2_1": {"name": "计算机存储", "keywords": ["ROM", "RAM", "Cache", "存储器"], "gespLevel": 2},
    "l2_2": {"name": "计算机网络基础", "keywords": ["网络", "TCP/IP", "IP地址", "协议"], "gespLevel": 2},
    "l2_3": {"name": "流程图", "keywords": ["流程图", "flowchart", "算法描述"], "gespLevel": 2},
    "l2_4": {"name": "ASCII编码", "keywords": ["ASCII", "编码", "字符编码"], "gespLevel": 2},
    "l2_5": {"name": "数据类型转换", "keywords": ["类型转换", "强制转换", "隐式转换"], "gespLevel": 2},
    "l2_6": {"name": "多层分支结构", "keywords": ["嵌套", "多层循环", "循环嵌套"], "gespLevel": 2},
    "l2_7": {"name": "数学函数", "keywords": ["abs", "sqrt", "max", "min", "rand", "数学函数"], "gespLevel": 2},
    "l2_8": {"name": "while和do-while", "keywords": ["while", "do-while", "条件循环"], "gespLevel": 2},
    "l2_9": {"name": "while深入", "keywords": ["while循环深入"], "gespLevel": 2},
    "l2_10": {"name": "数组进阶", "keywords": ["数组求最值", "数组查找", "数组统计"], "gespLevel": 2},
    "l2_11": {"name": "字符数组与字符串", "keywords": ["字符数组", "char数组", "字符串"], "gespLevel": 2},
    "l2_12": {"name": "string类型", "keywords": ["string", "string类", "getline"], "gespLevel": 2},
}

LESSON_TITLES = {1: "变量和数据的输入输出", 2: "算数运算符和字符型", 3: "比较运算符和 if 语句",
    4: "逻辑运算符和多分支语句", 5: "阶段复习与练习一", 6: "for 循环", 7: "for 循环进阶",
    8: "数组", 9: "数组进阶", 10: "阶段复习与练习二", 11: "while 循环",
    12: "格式化输入输出和浮点数", 13: "浮点数和数据类型转换", 14: "字符串和字符数组",
    15: "string 类型的字符串", 16: "阶段复习与练习三", 17: "函数基础", 18: "函数进阶",
    19: "结构体", 20: "循环的嵌套", 21: "枚举算法", 22: "模拟算法",
    23: "计算机基础与流程图", 24: "期末复习"}

REVIEW_LESSONS = {5: [1,2,3,4], 10: [6,7,8,9], 16: [11,12,13,14,15], 24: list(range(1,24))}

def extract_text(slide):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip())
        if shape.has_table:
            for row in shape.table.rows:
                row_text = [c.text.strip() for c in row.cells if c.text.strip()]
                if row_text:
                    texts.append(" | ".join(row_text))
    return "\n".join(texts)

def identify_points(text):
    found = set()
    text_lower = text.lower()
    for pid, pinfo in KNOWLEDGE_POINT_MAP.items():
        for kw in pinfo["keywords"]:
            if kw.lower() in text_lower:
                found.add(pid)
                break
    return list(found)

def get_lesson_num(filename):
    m = re.match(r'^(\d+)\.\s*', filename)
    if m: return int(m.group(1))
    m = re.match(r'^(\d+)\.pptx$', filename, re.I)
    if m: return int(m.group(1))
    return None

def analyze_file(filepath, lesson_num):
    filename = os.path.basename(filepath)
    title = LESSON_TITLES.get(lesson_num, f"Lesson {lesson_num}")
    is_review = lesson_num in REVIEW_LESSONS
    print(f"  Analyzing: {filename}")
    try:
        prs = Presentation(filepath)
    except Exception as e:
        print(f"    Error: {e}")
        return None
    pages = []
    all_points = set()
    for slide_idx, slide in enumerate(prs.slides, 1):
        content = extract_text(slide)
        if not content.strip(): continue
        points = identify_points(content)
        summary = content[:200].replace('\n', ' ').strip()
        if len(content) > 200: summary += "..."
        pages.append({"pageNum": slide_idx, "content": summary, "knowledgePoints": points})
        all_points.update(points)
    gesp_levels = set(KNOWLEDGE_POINT_MAP[pid]["gespLevel"] for pid in all_points if pid in KNOWLEDGE_POINT_MAP)
    return {"lessonNum": lesson_num, "title": title, "fileName": filename,
        "totalPages": len(prs.slides), "isReview": is_review, "pages": pages,
        "allKnowledgePoints": sorted(list(all_points)), "gespLevelsCovered": sorted(list(gesp_levels)),
        "reviewOfLessons": REVIEW_LESSONS.get(lesson_num) if is_review else None}

def main():
    print("=" * 60 + "\nLevel 1 PPT Analysis\n" + "=" * 60)
    ppt_files = [(get_lesson_num(f), f) for f in os.listdir(PPT_DIR) if f.lower().endswith('.pptx')]
    ppt_files = [(n, f) for n, f in ppt_files if n]
    ppt_files.sort(key=lambda x: x[0])
    print(f"\nFound {len(ppt_files)} PPT files\n")
    lessons = []
    all_kp = set()
    for lesson_num, filename in ppt_files:
        data = analyze_file(os.path.join(PPT_DIR, filename), lesson_num)
        if data:
            lessons.append(data)
            all_kp.update(data["allKnowledgePoints"])
    lessons.sort(key=lambda x: x["lessonNum"])
    for lesson in lessons:
        if lesson["isReview"] and lesson["reviewOfLessons"]:
            reviewed = set()
            for ln in lesson["reviewOfLessons"]:
                for rl in lessons:
                    if rl["lessonNum"] == ln:
                        reviewed.update(rl["allKnowledgePoints"])
                        break
            lesson["allKnowledgePoints"] = sorted(list(reviewed))
            lesson["gespLevelsCovered"] = sorted(list(set(
                KNOWLEDGE_POINT_MAP[p]["gespLevel"] for p in reviewed if p in KNOWLEDGE_POINT_MAP)))
    output = {"level": 1, "totalLessons": len(lessons), "lessons": lessons,
        "summary": {"totalKnowledgePoints": len(all_kp), "knowledgePointIds": sorted(list(all_kp)),
        "gespLevelsCovered": sorted(list(set(l for ls in lessons for l in ls["gespLevelsCovered"])))}}
    os.makedirs(os.path.dirname(OUTPUT_FILE), exist_ok=True)
    with open(OUTPUT_FILE, 'w', encoding='utf-8') as f:
        json.dump(output, f, ensure_ascii=False, indent=2)
    print("\n" + "=" * 60 + "\nAnalysis Complete!\n" + "=" * 60)
    print(f"Lessons: {len(lessons)}")
    print(f"Knowledge Points: {len(all_kp)}")
    print(f"GESP Levels: {output['summary']['gespLevelsCovered']}")
    print(f"Output: {OUTPUT_FILE}")

if __name__ == "__main__":
    main()
