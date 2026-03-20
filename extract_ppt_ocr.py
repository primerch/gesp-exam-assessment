#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
从PPT图片中提取文本（OCR）
"""

import os
import json
import io
from pptx import Presentation
from PIL import Image
import pytesseract

PPT_DIR = "/Users/reacher/Downloads/Archive/Level_4"
OUTPUT_DIR = "/Users/reacher/Developer/gesp/analysis_results"

def extract_text_from_ppt(filepath, max_slides=None):
    """使用OCR从PPT图片中提取文本"""
    prs = Presentation(filepath)
    
    results = []
    for i, slide in enumerate(prs.slides, 1):
        if max_slides and i > max_slides:
            break
        
        slide_text = []
        for shape in slide.shapes:
            # 尝试提取图片
            if hasattr(shape, 'image'):
                try:
                    image = shape.image
                    image_bytes = image.blob
                    pil_image = Image.open(io.BytesIO(image_bytes))
                    
                    # OCR识别中文和英文
                    text = pytesseract.image_to_string(pil_image, lang='chi_sim+eng')
                    if text.strip():
                        slide_text.append(text.strip())
                except Exception as e:
                    pass
            
            # 也尝试直接获取文本
            if hasattr(shape, 'text') and shape.text.strip():
                slide_text.append(shape.text.strip())
        
        if slide_text:
            results.append({
                "slide": i,
                "text": "\n".join(slide_text)
            })
    
    return results

def main():
    # 分析排列组合1作为测试
    files = [
        ("1-排列组合1.pptx", 1),
        ("2-排列组合2.pptx", 2),
    ]
    
    for filename, lesson_num in files:
        filepath = os.path.join(PPT_DIR, filename)
        print(f"OCR分析: {filename}")
        
        try:
            # 只分析前30页以节省时间
            slides = extract_text_from_ppt(filepath, max_slides=40)
            print(f"  提取了 {len(slides)} 页有内容的幻灯片")
            
            # 保存OCR结果
            result = {
                "lesson": lesson_num,
                "filename": filename,
                "slides": slides
            }
            
            output_file = os.path.join(OUTPUT_DIR, f"level4_lesson{lesson_num}_ocr.json")
            with open(output_file, 'w', encoding='utf-8') as f:
                json.dump(result, f, ensure_ascii=False, indent=2)
            print(f"  已保存: {output_file}")
            
            # 打印前3页内容样本
            for s in slides[:3]:
                print(f"\n  第 {s['slide']} 页:")
                print(f"  {s['text'][:300]}...")
                
        except Exception as e:
            print(f"  错误: {e}")

if __name__ == "__main__":
    main()
