#!/usr/bin/env python3
"""
Level 1 完整PPT分析脚本
逐页分析所有24个PPT文件，提取知识点并映射到GESP考点
"""

import json
import os
import re
from pptx import Presentation
from collections import defaultdict

# PPT文件目录
PPT_DIR = "/Users/reacher/Downloads/Archive/Level_1/"
OUTPUT_DIR = "/Users/reacher/Developer/gesp/analysis_results/"

# 确保输出目录存在
os.makedirs(OUTPUT_DIR, exist_ok=True)

# GESP 知识点映射定义 (Level 1 和 Level 2)
GESP_MAPPING = {
    # ========== GESP Level 1 考点 ==========
    "l1_1": {"name": "计算机基础", "desc": "计算机软硬件组成、操作系统基本概念"},
    "l1_2": {"name": "集成开发环境", "desc": "创建文件、编辑文件、保存文件、编译、解释、调试"},
    "l1_3": {"name": "程序基本框架", "desc": "#include、using namespace、main函数、return 0"},
    "l1_4": {"name": "输入输出语句", "desc": "cin/cout、scanf/printf的使用"},
    "l1_5": {"name": "变量定义", "desc": "标识符、关键字、常量、变量、命名规则、定义、初始化、赋值"},
    "l1_6": {"name": "基本数据类型", "desc": "int、long long、float、double、char、bool"},
    "l1_7": {"name": "算术运算", "desc": "+、-、*、/、%、整除"},
    "l1_8": {"name": "逻辑运算", "desc": "&&、||、!"},
    "l1_9": {"name": "关系运算", "desc": ">、>=、<、<=、==、!="},
    "l1_10": {"name": "三目运算", "desc": "条件?值1:值2"},
    "l1_11": {"name": "顺序结构", "desc": "程序按代码书写顺序依次执行"},
    "l1_12": {"name": "分支结构", "desc": "if、if-else、switch语句"},
    "l1_13": {"name": "循环结构", "desc": "for、while、do-while、break、continue"},
    "l1_14": {"name": "数组基础", "desc": "一维数组定义、初始化、访问、遍历"},
    
    # ========== GESP Level 2 考点 ==========
    "l2_1": {"name": "计算机存储", "desc": "ROM、RAM、Cache"},
    "l2_2": {"name": "计算机网络", "desc": "网络分类、TCP/IP模型、IP地址"},
    "l2_3": {"name": "流程图", "desc": "程序流程图符号和绘制方法"},
    "l2_4": {"name": "ASCII编码", "desc": "字符与ASCII码转换（空格32, 0=48, A=65, a=97）"},
    "l2_5": {"name": "数据类型转换", "desc": "强制转换、隐式转换"},
    "l2_6": {"name": "多层分支", "desc": "if嵌套、switch嵌套"},
    "l2_7": {"name": "多层循环", "desc": "循环嵌套"},
    "l2_8": {"name": "数学函数", "desc": "abs、sqrt、max、min、rand/srand"},
    "l2_9": {"name": "数组进阶", "desc": "数组求最值、查找、统计"},
    "l2_10": {"name": "字符数组与字符串", "desc": "字符数组存储字符串、字符串操作"},
    "l2_11": {"name": "string类型", "desc": "C++ string类型及成员函数"},
}

# 课程标题映射
LESSON_TITLES = {
    1: "变量和数据的输入输出",
    2: "算数运算符和字符型",
    3: "比较运算符和 if 语句",
    4: "逻辑运算符和多分支语句",
    5: "阶段复习与练习一",
    6: "for 循环",
    7: "for 循环进阶",
    8: "数组",
    9: "数组进阶",
    10: "阶段复习与练习二",
    11: "while 循环",
    12: "格式化输入输出和浮点数",
    13: "浮点数和数据类型转换",
    14: "字符串和字符数组",
    15: "string 类型的字符串",
    16: "阶段复习与练习三",
    17: "函数基础",
    18: "函数进阶",
    19: "结构体",
    20: "循环的嵌套",
    21: "枚举算法",
    22: "模拟算法",
    23: "计算机基础与流程图",
    24: "期末复习"
}

# PPT文件名映射
PPT_FILES = {
    1: "1. 变量和数据的输入输出.pptx",
    2: "2. 算数运算符和字符型.pptx",
    3: "3. 比较运算符和 if 语句.pptx",
    4: "4. 逻辑运算符和多分支语句.pptx",
    5: "5. 阶段复习与练习一.pptx",
    6: "6. for 循环.pptx",
    7: "7. for 循环进阶.pptx",
    8: "8. 数组.pptx",
    9: "9. 数组进阶.pptx",
    10: "10. 阶段复习与练习二.pptx",
    11: "11. while 循环.pptx",
    12: "12. 格式化输入输出和浮点数.pptx",
    13: "13. 浮点数和数据类型转换.pptx",
    14: "14. 字符串和字符数组.pptx",
    15: "15. string 类型的字符串.pptx",
    16: "16. 阶段复习与练习三.pptx",
    17: "17. 函数基础.pptx",
    18: "18. 函数进阶.pptx",
    19: "19. 结构体.pptx",
    20: "20. 循环的嵌套.pptx",
    21: "21.pptx",
    22: "22.pptx",
    23: "23.pptx",
    24: "24.pptx"
}


def extract_slide_content(slide, slide_num):
    """提取单页幻灯片的文本内容"""
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip())
        if shape.has_table:
            for row in shape.table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    texts.append(" | ".join(row_text))
    
    full_text = "\n".join(texts)
    return {
        "slide_number": slide_num,
        "content": texts,
        "full_text": full_text
    }


def identify_gesp_points(text):
    """识别文本中涉及的GESP知识点"""
    found_l1 = set()
    found_l2 = set()
    text_lower = text.lower()
    
    # ===== GESP Level 1 识别规则 =====
    # l1_1: 计算机基础
    if any(kw in text_lower for kw in ["计算机", "硬件", "软件", "操作系统", "cpu", "内存"]):
        found_l1.add("l1_1")
    
    # l1_2: 集成开发环境
    if any(kw in text_lower for kw in ["dev-c++", "编译", "调试", "保存", ".cpp", "ide"]):
        found_l1.add("l1_2")
    
    # l1_3: 程序基本框架
    if any(kw in text_lower for kw in ["#include", "using namespace", "main函数", "main(", "return 0"]):
        found_l1.add("l1_3")
    
    # l1_4: 输入输出语句
    if any(kw in text_lower for kw in ["cin", "cout", "scanf", "printf", "<<", ">>"]):
        found_l1.add("l1_4")
    
    # l1_5: 变量定义
    if any(kw in text_lower for kw in ["变量", "定义", "标识符", "关键字", "命名", "赋值", "初始化"]):
        found_l1.add("l1_5")
    
    # l1_6: 基本数据类型
    if any(kw in text_lower for kw in ["int", "long long", "float", "double", "char", "bool", "数据类型"]):
        found_l1.add("l1_6")
    
    # l1_7: 算术运算
    if any(kw in text_lower for kw in ["算术", "+", "-", "*", "/", "%", "取模", "整除", "++", "--"]):
        found_l1.add("l1_7")
    
    # l1_8: 逻辑运算
    if any(kw in text_lower for kw in ["逻辑", "&&", "||", "!", "并且", "或者", "非"]):
        found_l1.add("l1_8")
    
    # l1_9: 关系运算
    if any(kw in text_lower for kw in ["关系", "比较", ">", "<", "==", "!=", ">=", "<=", "等于", "大于", "小于"]):
        found_l1.add("l1_9")
    
    # l1_10: 三目运算
    if "?:" in text or "三目" in text:
        found_l1.add("l1_10")
    
    # l1_11: 顺序结构
    if "顺序" in text or "顺序结构" in text:
        found_l1.add("l1_11")
    
    # l1_12: 分支结构
    if any(kw in text_lower for kw in ["if", "else", "switch", "case", "分支", "条件", "判断"]):
        found_l1.add("l1_12")
    
    # l1_13: 循环结构
    if any(kw in text_lower for kw in ["for", "while", "do-while", "循环", "break", "continue"]):
        found_l1.add("l1_13")
    
    # l1_14: 数组基础
    if any(kw in text_lower for kw in ["数组", "一维数组", "arr[", "a[", "初始化", "遍历"]):
        found_l1.add("l1_14")
    
    # ===== GESP Level 2 识别规则 =====
    # l2_1: 计算机存储
    if any(kw in text_lower for kw in ["rom", "ram", "cache", "存储器"]):
        found_l2.add("l2_1")
    
    # l2_2: 计算机网络
    if any(kw in text_lower for kw in ["网络", "tcp/ip", "ip地址", "协议", "lan", "wan"]):
        found_l2.add("l2_2")
    
    # l2_3: 流程图
    if any(kw in text_lower for kw in ["流程图", "flowchart", "菱形", "矩形", "圆角矩形"]):
        found_l2.add("l2_3")
    
    # l2_4: ASCII编码
    if any(kw in text_lower for kw in ["ascii", "编码", "32", "48", "65", "97"]):
        found_l2.add("l2_4")
    
    # l2_5: 数据类型转换
    if any(kw in text_lower for kw in ["类型转换", "强制转换", "隐式转换", "(int)", "(double)", "(float)"]):
        found_l2.add("l2_5")
    
    # l2_6: 多层分支
    if any(kw in text_lower for kw in ["嵌套", "多层分支", "if嵌套", "switch嵌套"]):
        found_l2.add("l2_6")
    
    # l2_7: 多层循环
    if any(kw in text_lower for kw in ["循环嵌套", "双重循环", "多重循环", "双层for"]):
        found_l2.add("l2_7")
    
    # l2_8: 数学函数
    if any(kw in text_lower for kw in ["abs(", "sqrt(", "max(", "min(", "rand(", "srand(", "cmath"]):
        found_l2.add("l2_8")
    
    # l2_9: 数组进阶
    if any(kw in text_lower for kw in ["最值", "最大值", "最小值", "查找", "统计", "计数"]):
        found_l2.add("l2_9")
    
    # l2_10: 字符数组与字符串
    if any(kw in text_lower for kw in ["字符数组", "char[", "strlen", "strcpy", "字符型数组"]):
        found_l2.add("l2_10")
    
    # l2_11: string类型
    if any(kw in text_lower for kw in ["string", "#include <string>", ".length(", ".size(", ".substr(", ".find("]):
        found_l2.add("l2_11")
    
    return sorted(list(found_l1)), sorted(list(found_l2))


def extract_knowledge_points(slides_content, lesson_num):
    """提取课程的主要知识点"""
    knowledge = []
    all_text = "\n".join([s["full_text"] for s in slides_content])
    all_text_lower = all_text.lower()
    
    # 根据课程编号提取特定知识点
    if lesson_num == 1:
        if any(kw in all_text_lower for kw in ["#include", "main"]):
            knowledge.append({"name": "C++程序基本框架", "desc": "#include头文件、using namespace、main函数、return 0"})
        if "cout" in all_text_lower:
            knowledge.append({"name": "输出语句cout", "desc": "使用cout和<<进行标准输出"})
        if "cin" in all_text_lower:
            knowledge.append({"name": "输入语句cin", "desc": "使用cin和>>进行标准输入"})
        if "变量" in all_text:
            knowledge.append({"name": "变量定义与命名规则", "desc": "变量的定义、初始化、赋值、标识符命名规则"})
            
    elif lesson_num == 2:
        if any(kw in all_text for kw in ["算术", "+", "-", "*", "/"]):
            knowledge.append({"name": "算术运算符", "desc": "加减乘除、取模运算、整除"})
        if "char" in all_text_lower:
            knowledge.append({"name": "字符型数据类型", "desc": "char类型的定义和使用"})
        if "ascii" in all_text_lower:
            knowledge.append({"name": "ASCII编码", "desc": "字符与ASCII码的对应关系"})
            
    elif lesson_num == 3:
        if any(kw in all_text for kw in [">", "<", "=="]):
            knowledge.append({"name": "关系运算符", "desc": "大于、小于、等于、不等于等比较运算"})
        if "if" in all_text_lower:
            knowledge.append({"name": "if分支语句", "desc": "单分支if和双分支if-else语句"})
            
    elif lesson_num == 4:
        if any(kw in all_text_lower for kw in ["&&", "||", "!"]):
            knowledge.append({"name": "逻辑运算符", "desc": "逻辑与、逻辑或、逻辑非"})
        if "else if" in all_text_lower or "elseif" in all_text_lower:
            knowledge.append({"name": "多分支if-else if", "desc": "多条件分支判断"})
        if "switch" in all_text_lower:
            knowledge.append({"name": "switch语句", "desc": "switch-case-break-default结构"})
            
    elif lesson_num == 6:
        if "for" in all_text_lower:
            knowledge.append({"name": "for循环基础", "desc": "for循环语法：初始化、条件、更新"})
        if any(kw in all_text_lower for kw in ["break", "continue"]):
            knowledge.append({"name": "循环控制语句", "desc": "break跳出循环、continue跳过当前迭代"})
            
    elif lesson_num == 8:
        if "数组" in all_text:
            knowledge.append({"name": "一维数组定义", "desc": "数组的定义、初始化方法"})
            knowledge.append({"name": "数组元素访问", "desc": "通过下标访问数组元素"})
            
    elif lesson_num == 9:
        knowledge.append({"name": "数组求最值", "desc": "在数组中查找最大值和最小值"})
        knowledge.append({"name": "数组查找", "desc": "在数组中查找特定元素"})
        knowledge.append({"name": "数组统计", "desc": "数组元素求和、计数、平均值"})
        
    elif lesson_num == 11:
        if "while" in all_text_lower:
            knowledge.append({"name": "while循环", "desc": "当型循环的执行流程"})
        if "do-while" in all_text_lower or "do while" in all_text_lower:
            knowledge.append({"name": "do-while循环", "desc": "直到型循环的执行流程"})
            
    elif lesson_num == 12:
        if any(kw in all_text_lower for kw in ["float", "double"]):
            knowledge.append({"name": "浮点数类型", "desc": "float和double类型及精度"})
        if "printf" in all_text_lower:
            knowledge.append({"name": "printf格式化输出", "desc": "printf函数和格式控制符"})
        if "scanf" in all_text_lower:
            knowledge.append({"name": "scanf格式化输入", "desc": "scanf函数和格式控制符"})
            
    elif lesson_num == 13:
        knowledge.append({"name": "数据类型转换", "desc": "隐式转换和显式(强制)类型转换"})
        
    elif lesson_num == 14:
        knowledge.append({"name": "字符数组", "desc": "用字符数组存储字符串"})
        knowledge.append({"name": "字符串操作", "desc": "strlen、strcpy等字符串函数"})
        
    elif lesson_num == 15:
        knowledge.append({"name": "string类型", "desc": "C++ string类的使用"})
        knowledge.append({"name": "string成员函数", "desc": "length、substr、find等方法"})
        
    elif lesson_num == 17:
        knowledge.append({"name": "函数基础", "desc": "函数定义、声明、调用、参数、返回值"})
        
    elif lesson_num == 18:
        knowledge.append({"name": "数学函数", "desc": "abs、sqrt、max、min、rand/srand"})
        knowledge.append({"name": "函数参数传递", "desc": "值传递和引用传递"})
        
    elif lesson_num == 19:
        knowledge.append({"name": "结构体", "desc": "struct定义、成员访问"})
        
    elif lesson_num == 20:
        knowledge.append({"name": "循环嵌套", "desc": "多重循环的执行流程"})
        
    elif lesson_num in [21, 22, 23]:
        # 这些课可能是算法或综合复习
        knowledge.append({"name": "算法思维", "desc": "问题分析和算法设计"})
        
    elif lesson_num == 24:
        knowledge.append({"name": "期末综合复习", "desc": "全学期知识点回顾和巩固"})
    
    # 复习课特殊处理
    if lesson_num == 5:
        knowledge = [
            {"name": "阶段复习一", "desc": "复习第1-4课：变量、运算符、分支结构"}
        ]
    elif lesson_num == 10:
        knowledge = [
            {"name": "阶段复习二", "desc": "复习第6-9课：循环、数组"}
        ]
    elif lesson_num == 16:
        knowledge = [
            {"name": "阶段复习三", "desc": "复习第11-15课：while循环、浮点数、字符串"}
        ]
    
    return knowledge


def analyze_ppt_file(lesson_num):
    """分析单个PPT文件"""
    filename = PPT_FILES[lesson_num]
    filepath = os.path.join(PPT_DIR, filename)
    
    print(f"\n📄 分析第 {lesson_num} 课: {LESSON_TITLES[lesson_num]}")
    
    if not os.path.exists(filepath):
        print(f"  ❌ 文件不存在: {filepath}")
        return None
    
    try:
        prs = Presentation(filepath)
    except Exception as e:
        print(f"  ❌ 读取失败: {e}")
        return None
    
    # 逐页提取内容
    slides_content = []
    all_gesp_l1 = set()
    all_gesp_l2 = set()
    
    for slide_idx, slide in enumerate(prs.slides, 1):
        content = extract_slide_content(slide, slide_idx)
        if content["full_text"].strip():
            slides_content.append(content)
            
            # 识别GESP考点
            l1_points, l2_points = identify_gesp_points(content["full_text"])
            all_gesp_l1.update(l1_points)
            all_gesp_l2.update(l2_points)
    
    # 提取主要知识点
    knowledge_points = extract_knowledge_points(slides_content, lesson_num)
    
    # 构建结果
    result = {
        "level": 1,
        "lesson_number": lesson_num,
        "title": LESSON_TITLES[lesson_num],
        "file_name": filename,
        "total_slides": len(prs.slides),
        "slides_with_content": len(slides_content),
        "knowledge_points": knowledge_points,
        "gesp_level_1_points": sorted(list(all_gesp_l1)),
        "gesp_level_2_points": sorted(list(all_gesp_l2)),
        "slides_content": [
            {
                "slide_number": s["slide_number"],
                "content_preview": s["full_text"][:200] + "..." if len(s["full_text"]) > 200 else s["full_text"]
            }
            for s in slides_content[:10]  # 只保存前10页预览
        ]
    }
    
    # 保存单课分析结果
    output_file = os.path.join(OUTPUT_DIR, f"level1_lesson{lesson_num}.json")
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(result, f, ensure_ascii=False, indent=2)
    
    print(f"  ✅ 已保存: {output_file}")
    print(f"     总页数: {len(prs.slides)} | 有内容页: {len(slides_content)}")
    print(f"     GESP 1级考点: {len(all_gesp_l1)}个 | GESP 2级考点: {len(all_gesp_l2)}个")
    
    return result


def generate_summary(all_lessons):
    """生成Level 1总结文件"""
    print("\n" + "="*60)
    print("📊 生成 Level 1 总结报告...")
    print("="*60)
    
    # 汇总所有知识点
    all_l1_points = set()
    all_l2_points = set()
    all_knowledge = []
    
    for lesson in all_lessons:
        if lesson:
            all_l1_points.update(lesson.get("gesp_level_1_points", []))
            all_l2_points.update(lesson.get("gesp_level_2_points", []))
            all_knowledge.append({
                "lesson_number": lesson["lesson_number"],
                "title": lesson["title"],
                "knowledge_points": lesson.get("knowledge_points", []),
                "gesp_l1": lesson.get("gesp_level_1_points", []),
                "gesp_l2": lesson.get("gesp_level_2_points", [])
            })
    
    # 统计GESP 1级和2级考点覆盖情况
    total_l1 = 14  # GESP 1级共14个考点
    total_l2 = 11  # GESP 2级共11个考点
    covered_l1 = len(all_l1_points)
    covered_l2 = len(all_l2_points)
    
    # 计算缺失的考点
    all_l1_ids = set([f"l1_{i}" for i in range(1, 15)])
    all_l2_ids = set([f"l2_{i}" for i in range(1, 12)])
    missing_l1 = sorted(list(all_l1_ids - all_l1_points))
    missing_l2 = sorted(list(all_l2_ids - all_l2_points))
    
    summary = {
        "report_title": "Level 1 (24节课) GESP 考点分析报告",
        "generated_at": "2026-03-20",
        "total_lessons": 24,
        "lessons_analyzed": len([l for l in all_lessons if l]),
        
        # GESP 1级覆盖情况
        "gesp_level_1": {
            "total_points": total_l1,
            "covered_points": covered_l1,
            "coverage_rate": round(covered_l1 / total_l1 * 100, 1),
            "covered_ids": sorted(list(all_l1_points)),
            "missing_ids": missing_l1,
            "covered_details": [
                {"id": pid, "name": GESP_MAPPING[pid]["name"], "desc": GESP_MAPPING[pid]["desc"]}
                for pid in sorted(list(all_l1_points)) if pid in GESP_MAPPING
            ],
            "missing_details": [
                {"id": pid, "name": GESP_MAPPING[pid]["name"], "desc": GESP_MAPPING[pid]["desc"]}
                for pid in missing_l1 if pid in GESP_MAPPING
            ]
        },
        
        # GESP 2级覆盖情况
        "gesp_level_2": {
            "total_points": total_l2,
            "covered_points": covered_l2,
            "coverage_rate": round(covered_l2 / total_l2 * 100, 1),
            "covered_ids": sorted(list(all_l2_points)),
            "missing_ids": missing_l2,
            "covered_details": [
                {"id": pid, "name": GESP_MAPPING[pid]["name"], "desc": GESP_MAPPING[pid]["desc"]}
                for pid in sorted(list(all_l2_points)) if pid in GESP_MAPPING
            ],
            "missing_details": [
                {"id": pid, "name": GESP_MAPPING[pid]["name"], "desc": GESP_MAPPING[pid]["desc"]}
                for pid in missing_l2 if pid in GESP_MAPPING
            ]
        },
        
        # 每节课的知识点汇总
        "lesson_knowledge_summary": all_knowledge,
        
        # 考试建议
        "exam_recommendations": {
            "gesp_level_1": {
                "readiness": "ready" if covered_l1 >= 12 else "partial",
                "message": f"已覆盖 {covered_l1}/{total_l1} 个考点，建议完成全部课程后报考GESP 1级" if covered_l1 >= 12 else f"仅覆盖 {covered_l1}/{total_l1} 个考点，建议补充缺失知识点后再报考",
                "missing_critical": [GESP_MAPPING[pid]["name"] for pid in missing_l1[:5]]
            },
            "gesp_level_2": {
                "readiness": "ready" if covered_l2 >= 9 else "partial",
                "message": f"已覆盖 {covered_l2}/{total_l2} 个考点，建议完成全部课程后报考GESP 2级" if covered_l2 >= 9 else f"仅覆盖 {covered_l2}/{total_l2} 个考点，需要继续学习",
                "missing_critical": [GESP_MAPPING[pid]["name"] for pid in missing_l2[:5]]
            }
        }
    }
    
    # 保存总结文件
    summary_file = os.path.join(OUTPUT_DIR, "level1_summary.json")
    with open(summary_file, 'w', encoding='utf-8') as f:
        json.dump(summary, f, ensure_ascii=False, indent=2)
    
    print(f"\n✅ 总结文件已保存: {summary_file}")
    print(f"\n📈 Level 1 覆盖统计:")
    print(f"   GESP 1级: {covered_l1}/{total_l1} 个考点 ({summary['gesp_level_1']['coverage_rate']}%)")
    print(f"   GESP 2级: {covered_l2}/{total_l2} 个考点 ({summary['gesp_level_2']['coverage_rate']}%)")
    
    if missing_l1:
        print(f"\n⚠️  GESP 1级缺失考点 ({len(missing_l1)}个):")
        for pid in missing_l1[:5]:
            print(f"      - {GESP_MAPPING.get(pid, {}).get('name', pid)}")
    
    if missing_l2:
        print(f"\n⚠️  GESP 2级缺失考点 ({len(missing_l2)}个):")
        for pid in missing_l2[:5]:
            print(f"      - {GESP_MAPPING.get(pid, {}).get('name', pid)}")
    
    return summary


def main():
    """主函数"""
    print("="*60)
    print("🚀 Level 1 PPT 完整分析")
    print("="*60)
    print(f"输入目录: {PPT_DIR}")
    print(f"输出目录: {OUTPUT_DIR}")
    print("="*60)
    
    all_lessons = []
    
    # 分析所有24节课
    for lesson_num in range(1, 25):
        result = analyze_ppt_file(lesson_num)
        all_lessons.append(result)
    
    # 生成总结报告
    summary = generate_summary(all_lessons)
    
    print("\n" + "="*60)
    print("✅ 分析完成！")
    print("="*60)
    print(f"📁 输出文件:")
    print(f"   - {OUTPUT_DIR}level1_lesson1.json ~ level1_lesson24.json")
    print(f"   - {OUTPUT_DIR}level1_summary.json")
    print("="*60)


if __name__ == "__main__":
    main()
