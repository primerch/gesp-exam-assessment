#!/usr/bin/env python3
"""
分析 Level 1 所有 PPT 文件，提取知识点并映射到 GESP C++ 大纲
"""

from pptx import Presentation
import json
import os
import re
from pathlib import Path

# GESP 知识点映射定义
GESP_MAPPING = {
    # Level 1 考点
    "l1_1": "计算机基础：计算机软硬件组成、操作系统基本概念",
    "l1_2": "集成开发环境：创建文件、编辑文件、保存文件、编译、解释、调试",
    "l1_3": "程序框架：#include、using namespace、main函数、return 0",
    "l1_4": "输入输出：cin/cout、scanf/printf",
    "l1_5": "变量定义：标识符、关键字、常量、变量、命名规则、定义、初始化、赋值",
    "l1_6": "基本数据类型：int、long long、float、double、char、bool",
    "l1_7": "算术运算：+、-、*、/、%、整除",
    "l1_8": "逻辑运算：&&、||、!",
    "l1_9": "关系运算：>、>=、<、<=、==、!=",
    "l1_10": "三目运算：条件?值1:值2",
    "l1_11": "顺序结构",
    "l1_12": "分支结构：if、if-else、switch",
    "l1_13": "循环结构：for、while、do-while、break、continue",
    "l1_14": "数组基础：一维数组定义、初始化、访问、遍历",
    
    # Level 2 考点
    "l2_1": "计算机存储：ROM、RAM、Cache",
    "l2_2": "计算机网络：网络分类、TCP/IP模型、IP地址",
    "l2_3": "流程图：符号、绘制方法",
    "l2_4": "ASCII编码：字符与ASCII码转换（空格32, 0=48, A=65, a=97）",
    "l2_5": "数据类型转换：强制转换、隐式转换",
    "l2_6": "多层分支：if嵌套、switch嵌套",
    "l2_7": "多层循环：循环嵌套",
    "l2_8": "数学函数：abs、sqrt、max、min、rand/srand",
    "l2_9": "数组进阶：求最值、查找、统计",
    "l2_10": "字符数组与字符串",
    "l2_11": "string类型"
}

# 课程文件列表
LESSON_FILES = [
    ("1. 变量和数据的输入输出.pptx", "变量和数据的输入输出"),
    ("2. 算数运算符和字符型.pptx", "算数运算符和字符型"),
    ("3. 比较运算符和 if 语句.pptx", "比较运算符和 if 语句"),
    ("4. 逻辑运算符和多分支语句.pptx", "逻辑运算符和多分支语句"),
    ("5. 阶段复习与练习一.pptx", "阶段复习与练习一"),
    ("6. for 循环.pptx", "for 循环"),
    ("7. for 循环进阶.pptx", "for 循环进阶"),
    ("8. 数组.pptx", "数组"),
    ("9. 数组进阶.pptx", "数组进阶"),
    ("10. 阶段复习与练习二.pptx", "阶段复习与练习二"),
    ("11. while 循环.pptx", "while 循环"),
    ("12. 格式化输入输出和浮点数.pptx", "格式化输入输出和浮点数"),
    ("13. 浮点数和数据类型转换.pptx", "浮点数和数据类型转换"),
    ("14. 字符串和字符数组.pptx", "字符串和字符数组"),
    ("15. string 类型的字符串.pptx", "string 类型的字符串"),
    ("16. 阶段复习与练习三.pptx", "阶段复习与练习三"),
    ("17. 函数基础.pptx", "函数基础"),
    ("18. 函数进阶.pptx", "函数进阶"),
    ("19. 结构体.pptx", "结构体"),
    ("20. 循环的嵌套.pptx", "循环的嵌套"),
    ("21.pptx", None),  # 需要检查内容确定标题
    ("22.pptx", None),
    ("23.pptx", None),
    ("24.pptx", None),
]

def extract_text_from_slide(slide):
    """从幻灯片提取所有文本"""
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip())
        # 处理表格
        if hasattr(shape, "table"):
            for row in shape.table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    texts.append(" | ".join(row_text))
    return texts

def analyze_knowledge_points(content_texts, lesson_num):
    """分析文本内容，提取知识点"""
    all_text = "\n".join(content_texts).lower()
    knowledge_points = []
    gesp_l1 = set()
    gesp_l2 = set()
    
    # 根据课程内容匹配知识点
    # 这里需要根据实际的 PPT 内容进行详细分析
    
    return knowledge_points, list(gesp_l1), list(gesp_l2)

def analyze_lesson(lesson_num, filename, title_override=None):
    """分析单节课程"""
    ppt_path = f"/Users/reacher/Downloads/Archive/Level_1/{filename}"
    
    try:
        prs = Presentation(ppt_path)
    except Exception as e:
        print(f"Error loading {filename}: {e}")
        return None
    
    # 提取所有幻灯片内容
    slides_content = []
    all_texts = []
    
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = extract_text_from_slide(slide)
        slides_content.append({
            "slide_number": i,
            "content": slide_texts
        })
        all_texts.extend(slide_texts)
    
    # 尝试从第一页提取标题
    title = title_override
    if title is None and slides_content:
        # 从第一页提取标题
        first_slide_text = "\n".join(slides_content[0]["content"])
        # 通常标题在第一行
        lines = first_slide_text.strip().split('\n')
        if lines:
            title = lines[0].strip()
    
    result = {
        "level": 1,
        "lesson_number": lesson_num,
        "title": title or f"第{lesson_num}课",
        "file_name": filename,
        "total_slides": len(prs.slides),
        "slides_content": slides_content,
        "all_text": all_texts
    }
    
    return result

def save_raw_analysis(lesson_num, data):
    """保存原始分析结果"""
    output_path = f"/Users/reacher/Developer/gesp/analysis_results/level1_lesson{lesson_num}_raw.json"
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    print(f"Saved raw analysis: {output_path}")

def main():
    """主函数"""
    print("开始分析 Level 1 的所有 PPT 文件...")
    
    for lesson_num, (filename, title) in enumerate(LESSON_FILES, 1):
        print(f"\n分析第 {lesson_num} 课: {filename}")
        result = analyze_lesson(lesson_num, filename, title)
        if result:
            save_raw_analysis(lesson_num, result)
            print(f"  - 共 {result['total_slides']} 页幻灯片")
    
    print("\n所有课程原始内容已提取完成！")

if __name__ == "__main__":
    main()
