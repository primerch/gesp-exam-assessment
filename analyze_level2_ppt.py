#!/usr/bin/env python3
"""
深度分析 Level 2 (24节课) 的所有 PPT 内容
提取每节课知识点并映射到 GESP C++ 3-4级大纲
"""

import json
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

# GESP 3-4级知识点定义
GESP_LEVEL3_POINTS = {
    "l3_1": {"name": "数据编码", "category": "编码", "desc": "原码、反码、补码的概念", "keywords": ["原码", "反码", "补码", "编码", "负数表示"]},
    "l3_2": {"name": "进制转换", "category": "算法", "desc": "二进制、八进制、十进制、十六进制互转", "keywords": ["进制", "二进制", "八进制", "十进制", "十六进制", "转换"]},
    "l3_3": {"name": "位运算", "category": "运算", "desc": "与&、或|、非~、异或^、移位<<>>", "keywords": ["位运算", "与", "或", "非", "异或", "左移", "右移", "&", "|", "~", "^", "<<", ">>"]},
    "l3_4": {"name": "算法描述", "category": "算法", "desc": "自然语言、流程图、伪代码", "keywords": ["算法描述", "自然语言", "流程图", "伪代码"]},
    "l3_5": {"name": "枚举算法", "category": "算法", "desc": "穷举所有可能解逐一验证", "keywords": ["枚举", "穷举", "遍历", "所有可能"]},
    "l3_6": {"name": "模拟算法", "category": "算法", "desc": "按题目描述的过程一步步模拟", "keywords": ["模拟", "过程", "步骤"]},
    "l3_7": {"name": "递推算法", "category": "算法", "desc": "从已知条件出发推导结果", "keywords": ["递推", "推导", "迭代", " Fibonacci", "斐波那契"]},
    "l3_8": {"name": "函数基础", "category": "函数", "desc": "函数定义、调用、参数、返回值", "keywords": ["函数", "function", "调用", "参数", "返回值", "return"]},
}

GESP_LEVEL4_POINTS = {
    "l4_1": {"name": "指针基础", "category": "指针", "desc": "指针概念、定义、赋值、解引用", "keywords": ["指针", "pointer", "*", "&", "地址", "解引用", "内存"]},
    "l4_2": {"name": "二维数组", "category": "数据结构", "desc": "二维及多维数组的定义和使用", "keywords": ["二维数组", "多维数组", "[][]", "矩阵", "表格"]},
    "l4_3": {"name": "结构体", "category": "数据结构", "desc": "struct定义、结构体数组、结构体指针", "keywords": ["结构体", "struct", "自定义类型", "成员变量"]},
    "l4_4": {"name": "函数进阶", "category": "函数", "desc": "函数声明、形参实参、值传递、引用传递", "keywords": ["函数声明", "形参", "实参", "值传递", "引用传递", "传参"]},
    "l4_5": {"name": "变量作用域", "category": "函数", "desc": "全局变量与局部变量", "keywords": ["作用域", "全局变量", "局部变量", "生命周期"]},
    "l4_6": {"name": "排序算法", "category": "算法", "desc": "冒泡、选择、插入排序", "keywords": ["排序", "冒泡", "选择", "插入", "bubble", "selection", "insertion"]},
    "l4_7": {"name": "算法复杂度", "category": "算法", "desc": "时间复杂度、空间复杂度估算", "keywords": ["复杂度", "时间复杂度", "空间复杂度", "O(n)", "大O"]},
    "l4_8": {"name": "文件操作", "category": "文件", "desc": "文件读写、重定向", "keywords": ["文件", "file", "freopen", "fopen", "读写", "重定向"]},
    "l4_9": {"name": "高精度运算", "category": "算法", "desc": "大整数的加减乘除", "keywords": ["高精度", "大整数", "大数", "高精度加法", "高精度减法", "高精度乘法"]},
    "l4_10": {"name": "栈和队列", "category": "数据结构", "desc": "LIFO和FIFO的线性结构", "keywords": ["栈", "stack", "队列", "queue", "LIFO", "FIFO", "先进先出", "后进先出"]},
    "l4_11": {"name": "链表", "category": "数据结构", "desc": "单链表、双链表的基本操作", "keywords": ["链表", "linked list", "节点", "node", "next", "指针"]},
    "l4_12": {"name": "递归算法", "category": "算法", "desc": "函数自己调用自己的思想", "keywords": ["递归", "recursion", "递归调用", "递归函数", "递归出口"]},
    "l4_13": {"name": "二分查找", "category": "算法", "desc": "在有序序列中快速查找", "keywords": ["二分", "二分查找", "binary search", "折半", "对半"]},
    "l4_14": {"name": "贪心算法", "category": "算法", "desc": "局部最优推导全局最优", "keywords": ["贪心", "greedy", "局部最优", "最优解"]},
}

def extract_text_from_ppt(ppt_path):
    """从PPT中提取所有文本内容"""
    try:
        prs = Presentation(ppt_path)
        slides_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
            
            slides_content.append({
                "slide_num": slide_num,
                "text": "\n".join(slide_text)
            })
        
        return slides_content, len(prs.slides)
    except Exception as e:
        print(f"Error reading {ppt_path}: {e}")
        return [], 0

def identify_knowledge_points(text_content):
    """根据文本内容识别知识点"""
    text_lower = text_content.lower()
    
    found_points = {
        "level3": [],
        "level4": [],
        "general": []
    }
    
    # 检查GESP 3级知识点
    for point_id, point_info in GESP_LEVEL3_POINTS.items():
        for keyword in point_info["keywords"]:
            if keyword.lower() in text_lower:
                found_points["level3"].append({
                    "id": point_id,
                    "name": point_info["name"],
                    "category": point_info["category"]
                })
                break
    
    # 检查GESP 4级知识点
    for point_id, point_info in GESP_LEVEL4_POINTS.items():
        for keyword in point_info["keywords"]:
            if keyword.lower() in text_lower:
                found_points["level4"].append({
                    "id": point_id,
                    "name": point_info["name"],
                    "category": point_info["category"]
                })
                break
    
    return found_points

def get_lesson_title(slides_content):
    """从第一页提取课程标题"""
    if slides_content and len(slides_content) > 0:
        first_slide_text = slides_content[0]["text"]
        lines = [l.strip() for l in first_slide_text.split('\n') if l.strip()]
        for line in lines:
            if len(line) > 3 and len(line) < 100:
                return line
    return "Unknown"

def analyze_lesson(ppt_path, lesson_num):
    """分析单节课的PPT"""
    print(f"Analyzing L2-{lesson_num:02d}...")
    
    slides_content, total_slides = extract_text_from_ppt(ppt_path)
    
    if not slides_content:
        return None
    
    # 获取课程标题
    title = get_lesson_title(slides_content)
    
    # 分析所有页面的知识点
    all_knowledge = {
        "level3": [],
        "level4": [],
        "general": []
    }
    
    slide_analysis = []
    
    for slide in slides_content:
        points = identify_knowledge_points(slide["text"])
        slide_analysis.append({
            "slide_num": slide["slide_num"],
            "text_preview": slide["text"][:200] + "..." if len(slide["text"]) > 200 else slide["text"],
            "knowledge_points": points
        })
        
        # 合并知识点
        for p in points["level3"]:
            if p["id"] not in [x["id"] for x in all_knowledge["level3"]]:
                all_knowledge["level3"].append(p)
        for p in points["level4"]:
            if p["id"] not in [x["id"] for x in all_knowledge["level4"]]:
                all_knowledge["level4"].append(p)
    
    # 生成详细的知识点列表
    knowledge_points = []
    gesp_level_3_points = []
    gesp_level_4_points = []
    
    for p in all_knowledge["level3"]:
        knowledge_points.append(p["name"])
        gesp_level_3_points.append({
            "id": p["id"],
            "name": p["name"],
            "category": p["category"]
        })
    
    for p in all_knowledge["level4"]:
        knowledge_points.append(p["name"])
        gesp_level_4_points.append({
            "id": p["id"],
            "name": p["name"],
            "category": p["category"]
        })
    
    result = {
        "level": 2,
        "lesson_number": lesson_num,
        "title": title,
        "file_name": f"L2-{lesson_num:02d}.pptx",
        "total_slides": total_slides,
        "knowledge_points": knowledge_points,
        "gesp_level_3_points": gesp_level_3_points,
        "gesp_level_4_points": gesp_level_4_points,
        "slide_analysis": slide_analysis,
        "summary": f"本课共{total_slides}页，涵盖{len(knowledge_points)}个GESP相关知识点，" +
                   f"其中GESP 3级知识点{len(gesp_level_3_points)}个，GESP 4级知识点{len(gesp_level_4_points)}个。"
    }
    
    return result

def main():
    base_path = "/Users/reacher/Downloads/Archive/Level_2/"
    output_dir = "/Users/reacher/Developer/gesp/analysis_results/"
    
    # 确保输出目录存在
    os.makedirs(output_dir, exist_ok=True)
    
    # 分析所有24节课
    all_results = []
    
    for lesson_num in range(1, 25):
        ppt_filename = f"L2-{lesson_num:02d}.pptx"
        ppt_path = os.path.join(base_path, ppt_filename)
        
        if not os.path.exists(ppt_path):
            print(f"Warning: {ppt_path} not found")
            continue
        
        result = analyze_lesson(ppt_path, lesson_num)
        if result:
            all_results.append(result)
            
            # 保存单个课程分析结果
            output_file = os.path.join(output_dir, f"level2_lesson{lesson_num}.json")
            with open(output_file, 'w', encoding='utf-8') as f:
                json.dump(result, f, ensure_ascii=False, indent=2)
            print(f"  Saved: {output_file}")
    
    # 保存汇总结果
    summary_file = os.path.join(output_dir, "level2_all_lessons_summary.json")
    with open(summary_file, 'w', encoding='utf-8') as f:
        json.dump({
            "total_lessons": len(all_results),
            "lessons": [{"lesson": r["lesson_number"], "title": r["title"], 
                        "knowledge_count": len(r["knowledge_points"])} for r in all_results]
        }, f, ensure_ascii=False, indent=2)
    
    print(f"\nAnalysis complete!")
    print(f"Total lessons analyzed: {len(all_results)}")
    print(f"Results saved to: {output_dir}")

if __name__ == "__main__":
    main()
