#!/usr/bin/env python3
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt

def analyze_ppt(filepath):
    """分析PPT文件，提取每页的文本内容"""
    try:
        prs = Presentation(filepath)
        results = []
        
        for idx, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            results.append({
                'page': idx,
                'content': '\n'.join(slide_text)
            })
        
        return results
    except Exception as e:
        return [{'error': str(e)}]

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage: python analyze_ppt.py <pptx_file>")
        sys.exit(1)
    
    filepath = sys.argv[1]
    results = analyze_ppt(filepath)
    
    for slide in results:
        if 'error' in slide:
            print(f"Error: {slide['error']}")
        else:
            print(f"\n=== Page {slide['page']} ===")
            print(slide['content'][:2000])  # 限制每页输出长度
            print()
