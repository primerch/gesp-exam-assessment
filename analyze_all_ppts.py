#!/usr/bin/env python3
import os
import json
from pptx import Presentation

def analyze_ppt(filepath):
    """分析PPT文件，提取每页的文本内容"""
    try:
        prs = Presentation(filepath)
        pages = []
        
        for idx, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            content = '\n'.join(slide_text)
            pages.append({
                'page': idx,
                'content': content
            })
        
        return {
            'success': True,
            'total_pages': len(pages),
            'pages': pages
        }
    except Exception as e:
        return {
            'success': False,
            'error': str(e)
        }

def analyze_level(level_dir, level_name):
    """分析整个Level的所有PPT"""
    results = {
        'level': level_name,
        'lessons': []
    }
    
    # 获取所有pptx文件
    files = [f for f in os.listdir(level_dir) if f.endswith('.pptx')]
    files.sort()
    
    print(f"\n{'='*60}")
    print(f"分析 {level_name} - 共 {len(files)} 个课件")
    print(f"{'='*60}")
    
    for filename in files:
        filepath = os.path.join(level_dir, filename)
        print(f"\n分析: {filename}")
        
        analysis = analyze_ppt(filepath)
        
        if analysis['success']:
            lesson_info = {
                'filename': filename,
                'total_pages': analysis['total_pages'],
                'pages': analysis['pages']
            }
            results['lessons'].append(lesson_info)
            print(f"  ✓ 共 {analysis['total_pages']} 页")
        else:
            print(f"  ✗ 错误: {analysis.get('error', 'unknown')}")
    
    return results

if __name__ == '__main__':
    base_dir = "/Users/reacher/Downloads/Archive"
    
    # 分析所有Level
    all_results = {}
    
    for level_num in [1, 2, 3, 4]:
        level_dir = os.path.join(base_dir, f"Level_{level_num}")
        if os.path.exists(level_dir):
            result = analyze_level(level_dir, f"Level {level_num}")
            all_results[f"Level_{level_num}"] = result
    
    # 保存结果
    output_file = "/Users/reacher/Developer/gesp/ppt_analysis.json"
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(all_results, f, ensure_ascii=False, indent=2)
    
    print(f"\n{'='*60}")
    print(f"分析完成！结果已保存到: {output_file}")
    print(f"{'='*60}")
