
#!/usr/bin/env python3
"""
PPT内容提取脚本 - 提取每节课的知识点
"""

from pptx import Presentation
import json
import os
import sys
from pathlib import Path

def extract_ppt_content(ppt_path):
    """提取单个PPT的所有文本内容"""
    try:
        prs = Presentation(ppt_path)
        slides_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                slides_content.append({
                    "slide_num": slide_num,
                    "content": "\n".join(slide_text)
                })
        
        return slides_content
    except Exception as e:
        return [{"error": str(e)}]

def analyze_level(level_dir, output_file):
    """分析整个Level的所有PPT"""
    results = []
    
    ppt_files = sorted(Path(level_dir).glob("*.pptx"))
    
    for ppt_file in ppt_files:
        filename = ppt_file.name
        
        # 提取课序号
        lesson_num = None
        if filename.startswith('L') or filename[0].isdigit():
            # 尝试提取数字
            import re
            match = re.search(r'\d+', filename)
            if match:
                lesson_num = int(match.group())
        
        print(f"Processing: {filename}")
        
        content = extract_ppt_content(str(ppt_file))
        
        results.append({
            "filename": filename,
            "lesson_num": lesson_num,
            "title": filename.replace('.pptx', ''),
            "slides": content
        })
    
    # 保存结果
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(results, f, ensure_ascii=False, indent=2)
    
    print(f"Saved to: {output_file}")
    return results

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python extract_ppt.py <level_dir> <output_file>")
        sys.exit(1)
    
    level_dir = sys.argv[1]
    output_file = sys.argv[2]
    analyze_level(level_dir, output_file)
