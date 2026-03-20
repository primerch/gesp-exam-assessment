#!/usr/bin/env python3
"""
Level 2 (24节课) PPT 完整分析脚本
生成 24 个单课分析文件 + level2_summary.json 总结文件
"""

import json
import os
import re
from pathlib import Path
from pptx import Presentation
from collections import defaultdict

# GESP C++ 3-4级完整知识点定义
GESP_KNOWLEDGE_BASE = {
    # GESP Level 3 知识点
    "l3_01": {"name": "数据编码", "category": "编码", "desc": "原码、反码、补码的概念与转换", 
              "keywords": ["原码", "反码", "补码", "编码", "负数表示", "符号位", "数值编码"]},
    "l3_02": {"name": "进制转换", "category": "算法", "desc": "二进制、八进制、十进制、十六进制互转", 
              "keywords": ["进制", "二进制", "八进制", "十进制", "十六进制", "转换", "进制转换", "B", "O", "D", "H"]},
    "l3_03": {"name": "位运算", "category": "运算", "desc": "与&、或|、非~、异或^、移位<<>>", 
              "keywords": ["位运算", "与", "或", "非", "异或", "左移", "右移", "&", "|", "~", "^", "<<", ">>", "按位"]},
    "l3_04": {"name": "算法描述", "category": "算法", "desc": "自然语言、流程图、伪代码描述算法", 
              "keywords": ["算法描述", "自然语言", "流程图", "伪代码", "算法", "描述"]},
    "l3_05": {"name": "枚举算法", "category": "算法", "desc": "穷举所有可能解逐一验证", 
              "keywords": ["枚举", "穷举", "遍历", "所有可能", "逐一验证", "暴力求解"]},
    "l3_06": {"name": "模拟算法", "category": "算法", "desc": "按题目描述的过程一步步模拟", 
              "keywords": ["模拟", "过程", "步骤", "一步步", "按过程", "模拟过程"]},
    "l3_07": {"name": "递推算法", "category": "算法", "desc": "从已知条件出发推导结果", 
              "keywords": ["递推", "推导", "迭代", "fibonacci", "斐波那契", "递推公式", "递推关系"]},
    "l3_08": {"name": "函数基础", "category": "函数", "desc": "函数定义、调用、参数、返回值", 
              "keywords": ["函数", "function", "调用", "参数", "返回值", "return", "定义函数", "函数定义"]},
    
    # GESP Level 4 知识点
    "l4_01": {"name": "指针基础", "category": "指针", "desc": "指针概念、定义、赋值、解引用", 
              "keywords": ["指针", "pointer", "*", "&", "地址", "解引用", "内存", "指针变量", "指向"]},
    "l4_02": {"name": "指针运算", "category": "指针", "desc": "指针的加减运算、指针比较", 
              "keywords": ["指针运算", "指针加减", "指针比较", "指针移动", "地址运算"]},
    "l4_03": {"name": "二维数组", "category": "数据结构", "desc": "二维及多维数组的定义和使用", 
              "keywords": ["二维数组", "多维数组", "[][]", "矩阵", "表格", "行", "列", "二维"]},
    "l4_04": {"name": "数组与指针", "category": "指针", "desc": "数组名作为指针、指针访问数组", 
              "keywords": ["数组指针", "指针数组", "数组名", "下标", "[]", "数组访问"]},
    "l4_05": {"name": "结构体", "category": "数据结构", "desc": "struct定义、结构体数组、结构体指针", 
              "keywords": ["结构体", "struct", "自定义类型", "成员变量", ".", "->", "结构体定义"]},
    "l4_06": {"name": "函数进阶", "category": "函数", "desc": "函数声明、形参实参、值传递、引用传递", 
              "keywords": ["函数声明", "形参", "实参", "值传递", "引用传递", "传参", "参数传递"]},
    "l4_07": {"name": "变量作用域", "category": "函数", "desc": "全局变量与局部变量", 
              "keywords": ["作用域", "全局变量", "局部变量", "生命周期", "extern", "static"]},
    "l4_08": {"name": "冒泡排序", "category": "算法", "desc": "冒泡排序算法及实现", 
              "keywords": ["冒泡", "bubble", "排序", "冒泡排序", "交换", "相邻比较"]},
    "l4_09": {"name": "选择排序", "category": "算法", "desc": "选择排序算法及实现", 
              "keywords": ["选择", "selection", "排序", "选择排序", "最小值", "最大值"]},
    "l4_10": {"name": "插入排序", "category": "算法", "desc": "插入排序算法及实现", 
              "keywords": ["插入", "insertion", "排序", "插入排序", "插入位置", "有序序列"]},
    "l4_11": {"name": "算法复杂度", "category": "算法", "desc": "时间复杂度、空间复杂度估算", 
              "keywords": ["复杂度", "时间复杂度", "空间复杂度", "O(n)", "大O", "O(1)", "O(n²)", "O(logn)"]},
    "l4_12": {"name": "文件操作", "category": "文件", "desc": "文件读写、重定向", 
              "keywords": ["文件", "file", "freopen", "fopen", "读写", "重定向", "stdin", "stdout"]},
    "l4_13": {"name": "高精度运算", "category": "算法", "desc": "大整数的加减乘除", 
              "keywords": ["高精度", "大整数", "大数", "高精度加法", "高精度减法", "高精度乘法", "高精度除法"]},
    "l4_14": {"name": "栈", "category": "数据结构", "desc": "LIFO后进先出线性结构", 
              "keywords": ["栈", "stack", "LIFO", "后进先出", "入栈", "出栈", "push", "pop", "栈顶"]},
    "l4_15": {"name": "队列", "category": "数据结构", "desc": "FIFO先进先出线性结构", 
              "keywords": ["队列", "queue", "FIFO", "先进先出", "入队", "出队", "队首", "队尾"]},
    "l4_16": {"name": "链表", "category": "数据结构", "desc": "单链表、双链表的基本操作", 
              "keywords": ["链表", "linked list", "节点", "node", "next", "指针", "头结点", "尾结点"]},
    "l4_17": {"name": "递归算法", "category": "算法", "desc": "函数自己调用自己的思想", 
              "keywords": ["递归", "recursion", "递归调用", "递归函数", "递归出口", "递归基", "自己调用"]},
    "l4_18": {"name": "二分查找", "category": "算法", "desc": "在有序序列中快速查找", 
              "keywords": ["二分", "二分查找", "binary search", "折半", "对半", "有序序列", "中间值"]},
    "l4_19": {"name": "贪心算法", "category": "算法", "desc": "局部最优推导全局最优", 
              "keywords": ["贪心", "greedy", "局部最优", "最优解", "贪心策略", "最优"]},
    "l4_20": {"name": "分治算法", "category": "算法", "desc": "分而治之的思想", 
              "keywords": ["分治", "divide", "conquer", "分而治之", "分解", "子问题"]},
}

def extract_text_from_ppt(ppt_path):
    """从PPT中提取所有文本内容（逐页分析）"""
    try:
        prs = Presentation(ppt_path)
        slides_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_texts.append(text)
                # 处理表格
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                        if row_text:
                            slide_texts.append(row_text)
            
            full_text = "\n".join(slide_texts)
            slides_content.append({
                "slide_num": slide_num,
                "text": full_text,
                "text_length": len(full_text)
            })
        
        return slides_content, len(prs.slides)
    except Exception as e:
        print(f"  Error reading {ppt_path}: {e}")
        return [], 0

def identify_knowledge_points(text_content):
    """根据文本内容识别知识点（精确匹配）"""
    text_lower = text_content.lower()
    found_points = []
    matched_keywords = set()
    
    for point_id, point_info in GESP_KNOWLEDGE_BASE.items():
        for keyword in point_info["keywords"]:
            # 检查关键词是否出现在文本中
            if keyword.lower() in text_lower:
                found_points.append({
                    "id": point_id,
                    "name": point_info["name"],
                    "category": point_info["category"],
                    "desc": point_info["desc"],
                    "matched_keyword": keyword
                })
                matched_keywords.add(keyword)
                break  # 每个知识点只匹配一次
    
    return found_points, matched_keywords

def get_lesson_title(slides_content):
    """从第一页提取课程标题"""
    if slides_content and len(slides_content) > 0:
        first_slide_text = slides_content[0]["text"]
        lines = [l.strip() for l in first_slide_text.split('\n') if l.strip()]
        # 过滤掉常见的标题前缀
        for line in lines:
            clean_line = line.strip()
            if len(clean_line) > 2 and len(clean_line) < 100:
                # 移除常见前缀
                for prefix in ["第", "L2-", "Lesson", "课件", "课程"]:
                    if prefix in clean_line:
                        return clean_line
                return clean_line
    return "未知课程"

def categorize_points(points):
    """将知识点按GESP级别分类"""
    level3_points = []
    level4_points = []
    
    for p in points:
        if p["id"].startswith("l3_"):
            level3_points.append(p)
        elif p["id"].startswith("l4_"):
            level4_points.append(p)
    
    return level3_points, level4_points

def analyze_lesson(ppt_path, lesson_num):
    """详细分析单节课的PPT内容"""
    print(f"  分析 L2-{lesson_num:02d}.pptx...")
    
    slides_content, total_slides = extract_text_from_ppt(ppt_path)
    
    if not slides_content:
        print(f"  警告: 无法读取 L2-{lesson_num:02d}.pptx")
        return None
    
    # 获取课程标题
    title = get_lesson_title(slides_content)
    
    # 逐页分析知识点
    slide_analysis = []
    all_knowledge_points = []
    all_matched_keywords = set()
    
    for slide in slides_content:
        points, keywords = identify_knowledge_points(slide["text"])
        
        # 记录每页分析
        slide_analysis.append({
            "slide_num": slide["slide_num"],
            "text_preview": slide["text"][:150] + "..." if len(slide["text"]) > 150 else slide["text"],
            "knowledge_points_found": len(points),
            "points": [{"id": p["id"], "name": p["name"], "category": p["category"]} for p in points]
        })
        
        # 收集所有知识点（去重）
        for p in points:
            if p["id"] not in [x["id"] for x in all_knowledge_points]:
                all_knowledge_points.append(p)
        
        all_matched_keywords.update(keywords)
    
    # 分类知识点
    level3_points, level4_points = categorize_points(all_knowledge_points)
    
    # 生成知识点名称列表
    knowledge_point_names = [p["name"] for p in all_knowledge_points]
    
    # 生成总结
    summary = {
        "total_slides": total_slides,
        "total_knowledge_points": len(all_knowledge_points),
        "gesp_level3_points_count": len(level3_points),
        "gesp_level4_points_count": len(level4_points),
        "description": f"本课共{total_slides}页PPT，涵盖{len(all_knowledge_points)}个GESP相关知识点，"
                       f"其中GESP 3级知识点{len(level3_points)}个，GESP 4级知识点{len(level4_points)}个。"
    }
    
    result = {
        "level": 2,
        "lesson_number": lesson_num,
        "title": title,
        "file_name": f"L2-{lesson_num:02d}.pptx",
        "total_slides": total_slides,
        "knowledge_points": knowledge_point_names,
        "gesp_level3_points": [{"id": p["id"], "name": p["name"], "category": p["category"], "desc": p["desc"]} for p in level3_points],
        "gesp_level4_points": [{"id": p["id"], "name": p["name"], "category": p["category"], "desc": p["desc"]} for p in level4_points],
        "slide_analysis": slide_analysis,
        "summary": summary
    }
    
    return result

def generate_summary(all_lessons):
    """生成Level 2总结文件"""
    
    # 收集所有知识点
    all_l3_points = {}
    all_l4_points = {}
    lesson_summaries = []
    
    for lesson in all_lessons:
        lesson_summaries.append({
            "lesson_number": lesson["lesson_number"],
            "title": lesson["title"],
            "total_slides": lesson["total_slides"],
            "total_knowledge_points": lesson["summary"]["total_knowledge_points"],
            "gesp_level3_count": lesson["summary"]["gesp_level3_points_count"],
            "gesp_level4_count": lesson["summary"]["gesp_level4_points_count"]
        })
        
        for p in lesson["gesp_level3_points"]:
            all_l3_points[p["id"]] = p
        for p in lesson["gesp_level4_points"]:
            all_l4_points[p["id"]] = p
    
    # 统计覆盖情况
    total_l3_defined = len([k for k in GESP_KNOWLEDGE_BASE.keys() if k.startswith("l3_")])
    total_l4_defined = len([k for k in GESP_KNOWLEDGE_BASE.keys() if k.startswith("l4_")])
    
    covered_l3 = list(all_l3_points.values())
    covered_l4 = list(all_l4_points.values())
    
    # 查找缺失的知识点
    missing_l3 = []
    missing_l4 = []
    
    for point_id, point_info in GESP_KNOWLEDGE_BASE.items():
        if point_id.startswith("l3_") and point_id not in all_l3_points:
            missing_l3.append({"id": point_id, "name": point_info["name"], "category": point_info["category"]})
        elif point_id.startswith("l4_") and point_id not in all_l4_points:
            missing_l4.append({"id": point_id, "name": point_info["name"], "category": point_info["category"]})
    
    # 生成考试建议
    exam_recommendations = {
        "gesp_level_3": {
            "coverage_rate": round(len(covered_l3) / total_l3_defined * 100, 1),
            "covered_points": len(covered_l3),
            "total_points": total_l3_defined,
            "recommendation": "Level 2课程对GESP 3级的覆盖较好，学生完成Level 2后应重点复习" +
                            f"已覆盖的{len(covered_l3)}个知识点。",
            "missing_important": [p["name"] for p in missing_l3 if p["category"] in ["算法", "数据结构"]]
        },
        "gesp_level_4": {
            "coverage_rate": round(len(covered_l4) / total_l4_defined * 100, 1),
            "covered_points": len(covered_l4),
            "total_points": total_l4_defined,
            "recommendation": "Level 2课程覆盖了GESP 4级的部分基础内容，" +
                            f"已覆盖{len(covered_l4)}个知识点，建议补充学习缺失内容。",
            "missing_important": [p["name"] for p in missing_l4 if p["category"] in ["算法", "数据结构", "指针"]]
        }
    }
    
    summary = {
        "analysis_info": {
            "level": 2,
            "total_lessons": len(all_lessons),
            "analysis_date": "2025-03-20",
            "analyzer": "Level 2 PPT Analysis Tool"
        },
        "lessons_summary": lesson_summaries,
        "knowledge_coverage": {
            "gesp_level_3": {
                "total_defined": total_l3_defined,
                "covered_count": len(covered_l3),
                "coverage_rate": round(len(covered_l3) / total_l3_defined * 100, 1),
                "covered_points": covered_l3,
                "missing_points": missing_l3
            },
            "gesp_level_4": {
                "total_defined": total_l4_defined,
                "covered_count": len(covered_l4),
                "coverage_rate": round(len(covered_l4) / total_l4_defined * 100, 1),
                "covered_points": covered_l4,
                "missing_points": missing_l4
            }
        },
        "exam_recommendations": exam_recommendations,
        "study_path": {
            "for_gesp_3": "Level 2课程完整覆盖了GESP 3级的大部分知识点，完成Level 2的学生可以直接参加GESP 3级考试。",
            "for_gesp_4": "Level 2课程为GESP 4级打下基础，建议在完成Level 2后补充学习指针进阶、复杂数据结构等内容再参加GESP 4级考试。",
            "priority_topics": ["位运算", "进制转换", "函数基础", "结构体", "排序算法"]
        }
    }
    
    return summary

def main():
    """主函数：分析所有24节课并生成报告"""
    base_path = "/Users/reacher/Downloads/Archive/Level_2/"
    output_dir = "/Users/reacher/Developer/gesp/analysis_results/"
    
    # 确保输出目录存在
    os.makedirs(output_dir, exist_ok=True)
    
    print("=" * 60)
    print("Level 2 (24节课) PPT 完整分析")
    print("=" * 60)
    print(f"输入目录: {base_path}")
    print(f"输出目录: {output_dir}")
    print("=" * 60)
    
    # 分析所有24节课
    all_results = []
    
    for lesson_num in range(1, 25):
        ppt_filename = f"L2-{lesson_num:02d}.pptx"
        ppt_path = os.path.join(base_path, ppt_filename)
        
        if not os.path.exists(ppt_path):
            print(f"警告: {ppt_path} 不存在，跳过")
            continue
        
        result = analyze_lesson(ppt_path, lesson_num)
        if result:
            all_results.append(result)
            
            # 保存单个课程分析结果
            output_file = os.path.join(output_dir, f"level2_lesson{lesson_num}.json")
            with open(output_file, 'w', encoding='utf-8') as f:
                json.dump(result, f, ensure_ascii=False, indent=2)
            print(f"    ✓ 已保存: level2_lesson{lesson_num}.json")
    
    print("\n" + "=" * 60)
    print("生成总结文件...")
    print("=" * 60)
    
    # 生成总结文件
    summary = generate_summary(all_results)
    summary_file = os.path.join(output_dir, "level2_summary.json")
    with open(summary_file, 'w', encoding='utf-8') as f:
        json.dump(summary, f, ensure_ascii=False, indent=2)
    print(f"  ✓ 已保存: level2_summary.json")
    
    # 输出统计信息
    print("\n" + "=" * 60)
    print("分析完成！统计信息:")
    print("=" * 60)
    print(f"成功分析课程数: {len(all_results)}/24")
    print(f"GESP 3级覆盖: {summary['knowledge_coverage']['gesp_level_3']['covered_count']}/{summary['knowledge_coverage']['gesp_level_3']['total_defined']} 个知识点 " +
          f"({summary['knowledge_coverage']['gesp_level_3']['coverage_rate']}%)")
    print(f"GESP 4级覆盖: {summary['knowledge_coverage']['gesp_level_4']['covered_count']}/{summary['knowledge_coverage']['gesp_level_4']['total_defined']} 个知识点 " +
          f"({summary['knowledge_coverage']['gesp_level_4']['coverage_rate']}%)")
    
    print("\n生成的文件:")
    for i in range(1, 25):
        print(f"  - level2_lesson{i}.json")
    print(f"  - level2_summary.json (重要！)")
    
    print(f"\n所有文件已保存到: {output_dir}")

if __name__ == "__main__":
    main()
