#!/usr/bin/env python3
"""
Level 2 PPT 分批分析脚本
"""

import json
import os
import sys
from pptx import Presentation

# GESP C++ 3-4级知识点定义
GESP_KNOWLEDGE_BASE = {
    # Level 3
    "l3_01": {"name": "数据编码", "category": "编码", "keywords": ["原码", "反码", "补码", "编码", "负数表示"]},
    "l3_02": {"name": "进制转换", "category": "算法", "keywords": ["进制", "二进制", "八进制", "十进制", "十六进制", "转换"]},
    "l3_03": {"name": "位运算", "category": "运算", "keywords": ["位运算", "与", "或", "非", "异或", "左移", "右移", "&", "|", "~", "^"]},
    "l3_04": {"name": "算法描述", "category": "算法", "keywords": ["算法描述", "自然语言", "流程图", "伪代码"]},
    "l3_05": {"name": "枚举算法", "category": "算法", "keywords": ["枚举", "穷举", "遍历", "所有可能"]},
    "l3_06": {"name": "模拟算法", "category": "算法", "keywords": ["模拟", "过程", "步骤", "一步步"]},
    "l3_07": {"name": "递推算法", "category": "算法", "keywords": ["递推", "推导", "迭代", "斐波那契"]},
    "l3_08": {"name": "函数基础", "category": "函数", "keywords": ["函数", "调用", "参数", "返回值", "return"]},
    
    # Level 4
    "l4_01": {"name": "指针基础", "category": "指针", "keywords": ["指针", "pointer", "地址", "解引用", "内存"]},
    "l4_02": {"name": "指针运算", "category": "指针", "keywords": ["指针运算", "指针加减", "指针比较"]},
    "l4_03": {"name": "二维数组", "category": "数据结构", "keywords": ["二维数组", "多维数组", "矩阵", "行", "列"]},
    "l4_04": {"name": "数组与指针", "category": "指针", "keywords": ["数组指针", "指针数组", "数组名"]},
    "l4_05": {"name": "结构体", "category": "数据结构", "keywords": ["结构体", "struct", "成员变量"]},
    "l4_06": {"name": "函数进阶", "category": "函数", "keywords": ["函数声明", "形参", "实参", "值传递", "引用传递"]},
    "l4_07": {"name": "变量作用域", "category": "函数", "keywords": ["作用域", "全局变量", "局部变量"]},
    "l4_08": {"name": "冒泡排序", "category": "算法", "keywords": ["冒泡", "bubble", "排序"]},
    "l4_09": {"name": "选择排序", "category": "算法", "keywords": ["选择", "selection", "排序"]},
    "l4_10": {"name": "插入排序", "category": "算法", "keywords": ["插入", "insertion", "排序"]},
    "l4_11": {"name": "算法复杂度", "category": "算法", "keywords": ["复杂度", "时间复杂度", "空间复杂度", "O(n)"]},
    "l4_12": {"name": "文件操作", "category": "文件", "keywords": ["文件", "file", "freopen", "读写"]},
    "l4_13": {"name": "高精度运算", "category": "算法", "keywords": ["高精度", "大整数", "大数"]},
    "l4_14": {"name": "栈", "category": "数据结构", "keywords": ["栈", "stack", "LIFO", "后进先出"]},
    "l4_15": {"name": "队列", "category": "数据结构", "keywords": ["队列", "queue", "FIFO", "先进先出"]},
    "l4_16": {"name": "链表", "category": "数据结构", "keywords": ["链表", "linked list", "节点", "node"]},
    "l4_17": {"name": "递归算法", "category": "算法", "keywords": ["递归", "recursion", "递归调用"]},
    "l4_18": {"name": "二分查找", "category": "算法", "keywords": ["二分", "二分查找", "binary search"]},
    "l4_19": {"name": "贪心算法", "category": "算法", "keywords": ["贪心", "greedy", "局部最优"]},
    "l4_20": {"name": "分治算法", "category": "算法", "keywords": ["分治", "divide", "conquer"]},
}

def extract_and_analyze(ppt_path, lesson_num):
    """分析单节课"""
    prs = Presentation(ppt_path)
    
    all_text = []
    slide_analysis = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    slide_texts.append(text)
        full_text = "\n".join(slide_texts)
        all_text.append(full_text)
        slide_analysis.append({
            "slide_num": slide_num,
            "text_preview": full_text[:100] + "..." if len(full_text) > 100 else full_text
        })
    
    # 识别知识点
    full_content = "\n".join(all_text).lower()
    
    # 获取标题（第一页第一行）
    title = "未知课程"
    if all_text:
        first_lines = [l.strip() for l in all_text[0].split('\n') if l.strip()]
        if first_lines:
            title = first_lines[0][:50]
    
    found_points = []
    for point_id, point_info in GESP_KNOWLEDGE_BASE.items():
        for keyword in point_info["keywords"]:
            if keyword.lower() in full_content:
                found_points.append({
                    "id": point_id,
                    "name": point_info["name"],
                    "category": point_info["category"]
                })
                break
    
    # 去重
    seen_ids = set()
    unique_points = []
    for p in found_points:
        if p["id"] not in seen_ids:
            seen_ids.add(p["id"])
            unique_points.append(p)
    
    # 分类
    l3_points = [p for p in unique_points if p["id"].startswith("l3_")]
    l4_points = [p for p in unique_points if p["id"].startswith("l4_")]
    
    return {
        "level": 2,
        "lesson_number": lesson_num,
        "title": title,
        "file_name": f"L2-{lesson_num:02d}.pptx",
        "total_slides": len(prs.slides),
        "knowledge_points": [p["name"] for p in unique_points],
        "gesp_level3_points": l3_points,
        "gesp_level4_points": l4_points,
        "slide_analysis": slide_analysis[:5],  # 只保留前5页分析，减少文件大小
        "summary": {
            "total_slides": len(prs.slides),
            "total_knowledge_points": len(unique_points),
            "gesp_level3_points_count": len(l3_points),
            "gesp_level4_points_count": len(l4_points)
        }
    }

def main():
    base_path = "/Users/reacher/Downloads/Archive/Level_2/"
    output_dir = "/Users/reacher/Developer/gesp/analysis_results/"
    os.makedirs(output_dir, exist_ok=True)
    
    start_lesson = int(sys.argv[1]) if len(sys.argv) > 1 else 1
    end_lesson = int(sys.argv[2]) if len(sys.argv) > 2 else 24
    
    print(f"分析 Level 2 课程 {start_lesson} 到 {end_lesson}...")
    
    for lesson_num in range(start_lesson, end_lesson + 1):
        ppt_path = os.path.join(base_path, f"L2-{lesson_num:02d}.pptx")
        
        if not os.path.exists(ppt_path):
            print(f"  跳过: L2-{lesson_num:02d}.pptx 不存在")
            continue
        
        try:
            result = extract_and_analyze(ppt_path, lesson_num)
            output_file = os.path.join(output_dir, f"level2_lesson{lesson_num}.json")
            with open(output_file, 'w', encoding='utf-8') as f:
                json.dump(result, f, ensure_ascii=False, indent=2)
            print(f"  ✓ L2-{lesson_num:02d}: {result['title'][:30]}... " +
                  f"({result['summary']['total_slides']}页, " +
                  f"{result['summary']['total_knowledge_points']}个知识点)")
        except Exception as e:
            print(f"  ✗ L2-{lesson_num:02d}: 错误 - {e}")
    
    print("完成!")

if __name__ == "__main__":
    main()
