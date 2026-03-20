#!/usr/bin/env python3
"""
Level 4 PPT Analysis Script
Analyzes all 23 PPT files and extracts knowledge points mapping to GESP 7-8
"""

import os
import json
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from collections import defaultdict

# PPT Directory
PPT_DIR = "/Users/reacher/Downloads/Archive/Level_4/"
OUTPUT_FILE = "/Users/reacher/Developer/gesp/data/analysis/level4-analysis.json"

# GESP Level 7-8 Knowledge Point Mapping
GESP_MAPPING = {
    # Level 7 - Math Library
    "l7_1": ["三角函数", "sin", "cos", "tan", "math库"],
    "l7_2": ["对数函数", "log", "log10", "log2"],
    "l7_3": ["指数函数", "exp", "pow"],
    
    # Level 7 - Dynamic Programming
    "l7_4": ["二维动态规划", "二维DP", "二维状态"],
    "l7_5": ["DP优化", "动态规划优化", "最长上升子序列", "LIS", "最长公共子序列", "LCS", "区间DP", "滚动数组"],
    
    # Level 7 - Graph Theory
    "l7_6": ["图的定义", "图的基本概念", "邻接矩阵", "邻接表"],
    "l7_7": ["图的DFS", "图深度优先", "图遍历"],
    "l7_8": ["图的BFS", "图广度优先", "图层次遍历"],
    "l7_9": ["泛洪算法", "flood fill", "洪水填充", "连通块"],
    
    # Level 7 - Data Structures
    "l7_10": ["哈希表", "hash", "哈希", "散列表"],
    
    # Level 8 - Combinatorics
    "l8_1": ["计数原理", "加法原理", "乘法原理"],
    "l8_2": ["排列", "排列数", "全排列"],
    "l8_3": ["组合", "组合数", "杨辉三角"],
    
    # Level 8 - Algorithms
    "l8_4": ["倍增法", "快速幂", "ST表", "跳跃表"],
    "l8_5": ["代数基础", "方程", "一元一次", "二元一次"],
    "l8_6": ["平面几何", "面积", "周长", "几何"],
    
    # Level 8 - Graph Algorithms
    "l8_7": ["最小生成树", "MST", "kruskal", "prim"],
    "l8_8": ["最短路径", "dijkstra", "floyd", "单源最短", "多源最短"],
    "l8_9": ["图论综合", "图论应用"],
    
    # Level 8 - Analysis
    "l8_10": ["算法复杂度分析", "复杂度分析"],
    "l8_11": ["算法优化", "优化技巧", "时空优化"],
    "l8_12": ["区间动态规划", "区间DP", "石子合并", "括号匹配", "回文串"],
    "l8_13": ["搜索剪枝", "剪枝", "搜索优化", "记忆化搜索"],
}

# Lesson file mapping
LESSON_FILES = [
    ("1-排列组合1.pptx", 1, "排列组合1", False),
    ("2-排列组合2.pptx", 2, "排列组合2", False),
    ("3-质数筛法.pptx", 3, "质数筛法", False),
    ("4-区间动态规划1.pptx", 4, "区间动态规划1", False),
    ("5-区间动态规划2.pptx", 5, "区间动态规划2", False),
    ("6-阶段测试.pptx", 6, "阶段测试", True),
    ("7-参赛相关.pptx", 7, "参赛相关", False),
    ("8-枚举算法应用.pptx", 8, "枚举算法应用", False),
    ("9-贪心算法应用（复赛）.pptx", 9, "贪心算法应用（复赛）", False),
    ("10-模拟算法应用（复赛）.pptx", 10, "模拟算法应用（复赛）", False),
    ("11-搜索算法应用（复赛）.pptx", 11, "搜索算法应用（复赛）", False),
    ("12-复赛模拟比赛.pptx", 12, "复赛模拟比赛", True),
    ("13-计算机基础知识.pptx", 13, "计算机基础知识", False),
    ("14-数据结构1（初赛）.pptx", 14, "数据结构1（初赛）", False),
    ("15-数据结构2（初赛）.pptx", 15, "数据结构2（初赛）", False),
    ("16-数学专题1（初赛）.pptx", 16, "数学专题1（初赛）", False),
    ("17-数学专题2（初赛）.pptx", 17, "数学专题2（初赛）", False),
    ("18-字符串专题（初赛）.pptx", 18, "字符串专题（初赛）", False),
    ("19-二分算法专题（初赛）.pptx", 19, "二分算法专题（初赛）", False),
    ("20-排序专题（初赛）.pptx", 20, "排序专题（初赛）", False),
    ("21-递归专题（初赛）.pptx", 21, "递归专题（初赛）", False),
    ("22-动态规划（初赛）.pptx", 22, "动态规划（初赛）", False),
    ("23-初赛模拟比赛.pptx", 23, "初赛模拟比赛", True),
]

def extract_text_from_slide(slide):
    """Extract all text from a slide"""
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip())
    return "\n".join(texts)

def identify_knowledge_points(text):
    """Identify GESP knowledge points from text content"""
    points = set()
    text_lower = text.lower()
    
    for point_id, keywords in GESP_MAPPING.items():
        for keyword in keywords:
            if keyword.lower() in text_lower or keyword in text:
                points.add(point_id)
                break
    
    return list(points)

def analyze_ppt(filename, lesson_num, title, is_review):
    """Analyze a single PPT file"""
    filepath = os.path.join(PPT_DIR, filename)
    
    if not os.path.exists(filepath):
        print(f"Warning: {filename} not found")
        return None
    
    try:
        prs = Presentation(filepath)
    except Exception as e:
        print(f"Error opening {filename}: {e}")
        return None
    
    pages = []
    all_points = set()
    
    for i, slide in enumerate(prs.slides, 1):
        content = extract_text_from_slide(slide)
        points = identify_knowledge_points(content)
        
        page_data = {
            "pageNum": i,
            "content": content[:500] + "..." if len(content) > 500 else content,
            "knowledgePoints": points
        }
        pages.append(page_data)
        all_points.update(points)
    
    # Determine which lessons this is a review of
    review_of = None
    if lesson_num == 6:  # 阶段测试
        review_of = [1, 2, 3, 4, 5]
    elif lesson_num == 12:  # 复赛模拟比赛
        review_of = [7, 8, 9, 10, 11]
    elif lesson_num == 23:  # 初赛模拟比赛
        review_of = [13, 14, 15, 16, 17, 18, 19, 20, 21, 22]
    
    return {
        "lessonNum": lesson_num,
        "title": title,
        "fileName": filename,
        "totalPages": len(prs.slides),
        "isReview": is_review,
        "pages": pages,
        "allKnowledgePoints": sorted(list(all_points)),
        "gespLevelsCovered": list(set([7 if p.startswith("l7") else 8 for p in all_points])),
        "reviewOfLessons": review_of
    }

def main():
    print("Starting Level 4 PPT Analysis...")
    print(f"PPT Directory: {PPT_DIR}")
    print(f"Output File: {OUTPUT_FILE}")
    print()
    
    lessons = []
    
    for filename, lesson_num, title, is_review in LESSON_FILES:
        print(f"Analyzing: {filename} (Lesson {lesson_num})...")
        lesson_data = analyze_ppt(filename, lesson_num, title, is_review)
        if lesson_data:
            lessons.append(lesson_data)
            print(f"  - Pages: {lesson_data['totalPages']}")
            print(f"  - Knowledge Points: {len(lesson_data['allKnowledgePoints'])}")
            print(f"  - GESP Levels: {lesson_data['gespLevelsCovered']}")
            if lesson_data['reviewOfLessons']:
                print(f"  - Review of lessons: {lesson_data['reviewOfLessons']}")
    
    # Handle review lessons - aggregate knowledge points from previous lessons
    print("\nProcessing review lessons...")
    
    # Create lookup by lesson number
    lessons_by_num = {l['lessonNum']: l for l in lessons}
    
    for lesson in lessons:
        if lesson['isReview'] and lesson['reviewOfLessons']:
            all_review_points = set()
            for prev_lesson_num in lesson['reviewOfLessons']:
                if prev_lesson_num in lessons_by_num:
                    all_review_points.update(lessons_by_num[prev_lesson_num]['allKnowledgePoints'])
            
            # Update review lesson with aggregated points
            lesson['allKnowledgePoints'] = sorted(list(all_review_points))
            lesson['gespLevelsCovered'] = list(set([7 if p.startswith("l7") else 8 for p in all_review_points]))
            print(f"Lesson {lesson['lessonNum']} ({lesson['title']}): Aggregated {len(all_review_points)} knowledge points")
    
    # Build final output
    output = {
        "level": 4,
        "totalLessons": 23,
        "lessons": lessons
    }
    
    # Save to file
    with open(OUTPUT_FILE, 'w', encoding='utf-8') as f:
        json.dump(output, f, ensure_ascii=False, indent=2)
    
    print(f"\n✓ Analysis complete! Output saved to: {OUTPUT_FILE}")
    print(f"\nSummary:")
    print(f"  - Total lessons analyzed: {len(lessons)}")
    total_points = set()
    for lesson in lessons:
        total_points.update(lesson['allKnowledgePoints'])
    print(f"  - Total unique knowledge points: {len(total_points)}")
    print(f"  - GESP Level 7 points: {len([p for p in total_points if p.startswith('l7_')])}")
    print(f"  - GESP Level 8 points: {len([p for p in total_points if p.startswith('l8_')])}")

if __name__ == "__main__":
    main()
